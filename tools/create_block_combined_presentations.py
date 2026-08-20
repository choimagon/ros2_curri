#!/usr/bin/env python3
"""Create one teaching deck per Block by preserving every M-series deck in order."""
from __future__ import annotations

from copy import copy
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FONT = 'Noto Sans CJK KR'
MONO = 'DejaVu Sans Mono'
COURSE_TITLE = 'ROS 2 기반 AGV End-to-End 개발 커리큘럼'
COURSE_DECK_FILENAME = 'ROS_2_기반_AGV_End-to-End_개발_커리큘럼_전체_통합_따라하기.pptx'
NAVY, BLUE, GREEN, WHITE, LIGHT, GREY = (RGBColor(18, 44, 78), RGBColor(34, 91, 155), RGBColor(42, 139, 90), RGBColor(255, 255, 255), RGBColor(244, 247, 251), RGBColor(92, 102, 114))

BLOCKS = {
    'A': {'folder': 'blocks/A_ros2_basics', 'title': 'ROS 2 기초', 'modules': [('M01', 'ROS 2 node·topic 관찰'), ('M02', 'workspace와 Python·C++ 패키지'), ('M03', 'Publisher / Subscriber'), ('M04', 'TF2 좌표계')]},
    'B': {'folder': 'blocks/B_robot_build', 'title': 'AGV 로봇 제작', 'modules': [('M05', '단일 URDF 기본 모델'), ('M06', 'Xacro macro·property'), ('M07', '물리·collision·inertia'), ('M08', 'World와 spawn')]},
    'C': {'folder': 'blocks/C_drive_visualization', 'title': '주행과 시각화', 'modules': [('M09', 'Differential Drive와 odom'), ('M10', 'ros_gz_bridge'), ('M11', 'RViz 통합 검증')]},
    'D': {'folder': 'blocks/D_sensors', 'title': '센서', 'modules': [('M12', 'Camera'), ('M13', '2D LiDAR'), ('M14', 'IMU'), ('M15', 'QoS·sim time·운영')]},
    'E': {'folder': 'blocks/E_autonomy_logic', 'title': '제어·인지·미션', 'modules': [('M16', 'ros2_control'), ('M17', 'Camera·YOLO / fallback'), ('M18', 'LiDAR 안전 판단'), ('M19', 'Mission FSM'), ('M20', 'PID 주행 제어')]},
    'F': {'folder': 'blocks/F_integration', 'title': '통합·재현·최종 프로젝트', 'modules': [('M21', 'YAML·Launch·rosbag2'), ('M22', '자율 AGV 최종 미션')]},
}


def set_text(shape, text: str, size: int, color: RGBColor, bold: bool = False, font: str = FONT, align=PP_ALIGN.LEFT) -> None:
    frame = shape.text_frame
    frame.clear(); frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.08)
    frame.margin_top = frame.margin_bottom = Inches(0.04)
    paragraph = frame.paragraphs[0]
    paragraph.text = text; paragraph.alignment = align
    paragraph.font.name = font; paragraph.font.size = Pt(size); paragraph.font.bold = bold; paragraph.font.color.rgb = color


def text_box(slide, x, y, w, h, text: str, size: int, color: RGBColor, bold: bool = False, font: str = FONT, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shape, text, size, color, bold, font, align)
    return shape


def box(slide, x, y, w, h, fill: RGBColor, line: RGBColor, rounded=False):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    return shape


def copy_font(source, target) -> None:
    target.name = source.name
    target.size = source.size
    target.bold = source.bold
    target.italic = source.italic
    target.underline = source.underline
    try:
        target.color.rgb = source.color.rgb
    except (AttributeError, TypeError):
        pass


def copy_text_frame(source, target) -> None:
    target.clear(); target.word_wrap = source.word_wrap; target.vertical_anchor = source.vertical_anchor
    target.margin_left, target.margin_right = source.margin_left, source.margin_right
    target.margin_top, target.margin_bottom = source.margin_top, source.margin_bottom
    for index, paragraph in enumerate(source.paragraphs):
        destination = target.paragraphs[0] if index == 0 else target.add_paragraph()
        destination.alignment = paragraph.alignment; destination.level = paragraph.level
        destination.space_after = paragraph.space_after; destination.space_before = paragraph.space_before
        copy_font(paragraph.font, destination.font)
        for run in paragraph.runs:
            copied = destination.add_run(); copied.text = run.text; copy_font(run.font, copied.font)


def copy_fill(source, target) -> None:
    try:
        if source.fill.type is not None:
            target.fill.solid(); target.fill.fore_color.rgb = source.fill.fore_color.rgb
    except (AttributeError, TypeError):
        pass
    try:
        target.line.color.rgb = source.line.color.rgb; target.line.width = source.line.width
    except (AttributeError, TypeError):
        pass


def copy_slide(source, presentation) -> None:
    destination = presentation.slides.add_slide(presentation.slide_layouts[6])
    for shape in source.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            destination.shapes.add_picture(BytesIO(shape.image.blob), shape.left, shape.top, shape.width, shape.height)
        elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
            copied = destination.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)
            copy_fill(shape, copied)
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            copied = destination.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
            copy_text_frame(shape.text_frame, copied.text_frame)
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            copied = destination.shapes.add_shape(shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height)
            copy_fill(shape, copied)
            if shape.has_text_frame:
                copy_text_frame(shape.text_frame, copied.text_frame)
        else:
            # The M decks use text, auto-shapes, pictures, and straight connectors.
            # Retain any future unsupported shape as a short visible warning.
            text_box(destination, 0.5, 0.5, 12.2, 0.3, f'[복사하지 않은 도형: {shape.shape_type}]', 10, GREY)
    try:
        destination.notes_slide.notes_text_frame.text = source.notes_slide.notes_text_frame.text
    except AttributeError:
        pass


def separator(presentation, block_id: str, module_id: str, module_title: str, position: int, total: int) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box(slide, 0, 0, 13.333, 1.05, NAVY, NAVY)
    text_box(slide, 0.7, 0.30, 12.0, 0.34, f'Block {block_id} · {position}/{total} · {module_id}', 20, WHITE, True, align=PP_ALIGN.CENTER)
    text_box(slide, 0.9, 2.2, 11.5, 0.7, module_title, 34, NAVY, True, align=PP_ALIGN.CENTER)
    panel = box(slide, 1.05, 3.4, 11.25, 1.45, LIGHT, LIGHT, rounded=True)
    set_text(panel, f'{module_id}의 독립 PPT 전체를 바로 다음 슬라이드부터 순서대로 진행합니다.\n이 Block 통합본에서도 Starter → 파일 작성 → 실행 → 검증 → 오류 대응 → Checkpoint 순서는 바뀌지 않습니다.', 19, NAVY, True, align=PP_ALIGN.CENTER)
    text_box(slide, 1.0, 6.45, 11.4, 0.32, '모듈을 건너뛰지 말고, 이전 M의 Complete가 준비된 뒤 다음 M으로 이동합니다.', 15, GREEN, True, align=PP_ALIGN.CENTER)


def add_summary(presentation, block_id: str, metadata: dict) -> None:
    modules = metadata['modules']
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box(slide, 0, 0, 13.333, 1.15, NAVY, NAVY)
    text_box(slide, 0.7, 0.22, 12.0, 0.28, COURSE_TITLE, 13, WHITE, True, align=PP_ALIGN.CENTER)
    text_box(slide, 0.7, 0.58, 12.0, 0.28, f'Block {block_id} · M 시리즈 통합 따라 하기', 19, WHITE, True, align=PP_ALIGN.CENTER)
    text_box(slide, 0.7, 2.0, 12.0, 0.76, metadata['title'], 36, NAVY, True, align=PP_ALIGN.CENTER)
    panel = box(slide, 1.25, 3.25, 10.85, 1.5, LIGHT, LIGHT, rounded=True)
    set_text(panel, '이 파일은 Block 안의 모든 M 시리즈 PPT를 순서와 발표자 노트까지 유지해 한 번에 볼 수 있도록 합친 자료입니다.\n개별 M PPT는 실습 때, 이 통합본은 Block 전체 흐름을 복습·강의할 때 사용합니다.', 19, NAVY, True, align=PP_ALIGN.CENTER)
    text_box(slide, 0.9, 6.5, 11.7, 0.28, f'포함 모듈: {" → ".join(module_id for module_id, _ in modules)}', 16, GREEN, True, align=PP_ALIGN.CENTER)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box(slide, 0.75, 0.55, 12.0, 0.5, f'Block {block_id} 커리큘럼 요약', 30, NAVY, True, align=PP_ALIGN.CENTER)
    for index, (module_id, description) in enumerate(modules):
        y = 1.35 + index * (4.75 / len(modules))
        panel = box(slide, 0.85, y, 11.65, (4.20 / len(modules)), WHITE, BLUE, rounded=True)
        text_box(slide, 1.12, y + 0.12, 1.25, 0.27, module_id, 19, BLUE, True)
        text_box(slide, 2.45, y + 0.10, 9.4, 0.30, description, 18, NAVY, True)
    text_box(slide, 0.9, 6.62, 11.5, 0.26, '각 모듈의 완료 조건을 통과한 뒤에만 다음 모듈로 진행합니다.', 14, GREEN, True, align=PP_ALIGN.CENTER)


def build_block(block_id: str, metadata: dict) -> Path:
    folder = ROOT / metadata['folder']
    presentation = Presentation()
    presentation.slide_width, presentation.slide_height = Inches(13.333333), Inches(7.5)
    add_summary(presentation, block_id, metadata)
    for index, (module_id, module_title) in enumerate(metadata['modules'], start=1):
        module_dir = next(folder.glob(f'{module_id}_*'))
        module_deck = next(module_dir.glob(f'{module_id}_*.pptx'))
        separator(presentation, block_id, module_id, module_title, index, len(metadata['modules']))
        source = Presentation(module_deck)
        for source_slide in source.slides:
            copy_slide(source_slide, presentation)
    output = folder / f'Block_{block_id}_M시리즈_통합_따라하기.pptx'
    presentation.save(output)
    return output


def add_course_cover(presentation) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box(slide, 0, 0, 13.333, 1.18, NAVY, NAVY)
    text_box(slide, 0.70, 0.34, 12.0, 0.36, '전체 통합 강의·복습 자료', 21, WHITE, True, align=PP_ALIGN.CENTER)
    text_box(slide, 0.75, 1.95, 11.85, 0.85, COURSE_TITLE, 35, NAVY, True, align=PP_ALIGN.CENTER)
    panel = box(slide, 1.12, 3.35, 11.10, 1.54, LIGHT, LIGHT, rounded=True)
    set_text(panel, 'Block A → B → C → D → E → F\nM01부터 M22까지의 개별 실습 PPT를 순서와 발표자 노트까지 유지해 한 파일로 합쳤습니다.', 20, NAVY, True, align=PP_ALIGN.CENTER)
    text_box(slide, 0.85, 6.50, 11.65, 0.30, '개별 M PPT는 직접 실습에, 이 전체 통합본은 전체 흐름 강의·복습에 사용합니다.', 15, GREEN, True, align=PP_ALIGN.CENTER)


def add_course_block_separator(presentation, block_id: str, metadata: dict, position: int) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box(slide, 0, 0, 13.333, 1.12, NAVY, NAVY)
    text_box(slide, 0.7, 0.31, 12.0, 0.34, f'{COURSE_TITLE} · Block {position}/6', 18, WHITE, True, align=PP_ALIGN.CENTER)
    text_box(slide, 0.9, 2.15, 11.55, 0.65, f'Block {block_id} — {metadata["title"]}', 34, NAVY, True, align=PP_ALIGN.CENTER)
    module_text = ' → '.join(f'{module_id} {description}' for module_id, description in metadata['modules'])
    panel = box(slide, 1.0, 3.52, 11.35, 1.35, LIGHT, LIGHT, rounded=True)
    set_text(panel, f'이 Block에서 진행할 모듈\n{module_text}', 17, NAVY, True, align=PP_ALIGN.CENTER)
    text_box(slide, 0.8, 6.47, 11.75, 0.31, '아래부터 각 M의 독립 PPT 전체가 시작됩니다. 이전 M의 checkpoint를 통과한 뒤 다음 M으로 이동합니다.', 14, GREEN, True, align=PP_ALIGN.CENTER)


def build_course() -> Path:
    """Combine every M01–M22 deck once, grouped by Block A–F."""
    presentation = Presentation()
    presentation.slide_width, presentation.slide_height = Inches(13.333333), Inches(7.5)
    add_course_cover(presentation)
    for block_position, (block_id, metadata) in enumerate(BLOCKS.items(), start=1):
        folder = ROOT / metadata['folder']
        add_course_block_separator(presentation, block_id, metadata, block_position)
        for module_position, (module_id, module_title) in enumerate(metadata['modules'], start=1):
            module_dir = next(folder.glob(f'{module_id}_*'))
            module_deck = next(module_dir.glob(f'{module_id}_*.pptx'))
            separator(presentation, block_id, module_id, module_title, module_position, len(metadata['modules']))
            source = Presentation(module_deck)
            for source_slide in source.slides:
                copy_slide(source_slide, presentation)
    output = ROOT / COURSE_DECK_FILENAME
    presentation.save(output)
    return output


if __name__ == '__main__':
    for identifier, definition in BLOCKS.items():
        output = build_block(identifier, definition)
        print(output.relative_to(ROOT))
    print(build_course().relative_to(ROOT))
