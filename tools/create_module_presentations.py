#!/usr/bin/env python3
"""Build the M01--M22 follow-along teaching decks and delivery artifacts.

The content and layout follow ``ROS2_Gazebo_AGV_따라하기형_PPT_제작_가이드.docx``:
every module has a starting state, one-action code slides with complete files,
build/run instructions, an actual validation capture, objective checks, common
errors, a checkpoint, and presenter notes.
"""
from __future__ import annotations

import hashlib
import ast
import os
import shutil
import subprocess
import textwrap
from pathlib import Path
from xml.dom import minidom

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
WORKSPACE = ROOT / "agv_ws"
SRC = WORKSPACE / "src"
FONT = "Noto Sans CJK KR"
MONO = "DejaVu Sans Mono"
FONT_FILE = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"

NAVY = RGBColor(18, 44, 78)
BLUE = RGBColor(34, 91, 155)
GREEN = RGBColor(42, 139, 90)
ORANGE = RGBColor(219, 126, 34)
RED = RGBColor(190, 59, 59)
GREY = RGBColor(92, 102, 114)
LIGHT = RGBColor(244, 247, 251)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(29, 35, 45)


# source paths are relative to ~/ros2_curri/agv_ws/src.  The files are read directly into
# the deck, so the slides never depend on an unlisted external snippet.
MODULES = [
    {
        "id": "M01", "block": "A", "folder": "blocks/A_ros2_basics/M01_ros2_concepts",
        "filename": "M01_ROS2_기초와_분산_로봇_소프트웨어.pptx",
        "title": "ROS 2 기초와 분산 로봇 소프트웨어",
        "goal": "여러 ROS 2 노드가 Topic으로 메시지를 주고받는 구조를 직접 관찰한다.",
        "completion": "talker/listener와 node·topic·echo·info 명령으로 연결을 확인한다.",
        "start": "새 Ubuntu 터미널. 아직 workspace나 사용자 코드는 필요 없다.",
        "previous": "과정 시작", "next": "M02 workspace 생성",
        "flow": ["talker\nNode", "/chatter\nTopic", "listener\nNode"],
        "files": [],
        "actions": [
            ("ROS 2 환경을 등록한다", "새 터미널은 ROS 명령을 아직 모른다.", "source /opt/ros/jazzy/setup.bash\nros2 --help | head"),
            ("두 노드를 서로 다른 터미널에서 실행한다", "Publisher와 Subscriber는 동시에 살아 있어야 한다.", "# 터미널 1\nros2 run demo_nodes_cpp talker\n\n# 새 터미널 2\nsource /opt/ros/jazzy/setup.bash\nros2 run demo_nodes_py listener"),
            ("연결을 명령으로 관찰한다", "코드 밖에서도 이름·타입·연결 수를 확인할 수 있다.", "ros2 node list\nros2 topic list\nros2 topic info /chatter\nros2 topic echo /chatter --once"),
        ],
        "run": "source /opt/ros/jazzy/setup.bash\nros2 run demo_nodes_cpp talker\n# 새 터미널: ros2 run demo_nodes_py listener",
        "validate": "ros2 node list\nros2 topic info /chatter\nros2 topic echo /chatter --once",
        "capture": "ros2 pkg prefix demo_nodes_cpp; ros2 interface show std_msgs/msg/String",
        "errors": [("ros2: command not found", "source /opt/ros/jazzy/setup.bash를 실행한다."), ("echo가 출력되지 않음", "talker 실행·/chatter 이름·publisher 수를 순서대로 확인한다.")],
        "mini": "listener를 Ctrl-C로 멈춘 뒤 topic info의 subscriber 수가 어떻게 바뀌는지 확인한다.",
    },
    {
        "id": "M02", "block": "A", "folder": "blocks/A_ros2_basics/M02_workspace_package",
        "filename": "M02_Workspace_생성과_빌드.pptx", "title": "Workspace 생성과 Python·C++ 패키지 빌드",
        "goal": "~/ros2_curri/agv_ws를 만들고 Python ament_python과 C++ ament_cmake 패키지를 빌드한다.",
        "completion": "build/install/log와 두 언어의 실행 파일을 source 뒤에 확인한다.",
        "start": "M01을 끝낸 터미널. ~/ros2_curri/agv_ws가 없거나, 기존 폴더를 별도로 백업한 상태.",
        "previous": "M01 ROS 2 source", "next": "M03 Python Pub/Sub",
        "flow": ["src\n소스", "colcon build", "install\noverlay"],
        "files": ["agv_control/package.xml", "agv_control/setup.py", "agv_cpp_examples/CMakeLists.txt", "agv_cpp_examples/src/status_publisher.cpp"],
        "actions": [("workspace 뼈대를 만든다", "colcon은 src 아래 패키지를 찾아 build/install/log를 만든다.", "mkdir -p ~/ros2_curri/agv_ws/src\ncd ~/ros2_curri/agv_ws\ncolcon build --symlink-install"), ("새 터미널에서 overlay를 source한다", "기본 ROS와 내가 만든 패키지는 서로 다른 setup 파일이다.", "source /opt/ros/jazzy/setup.bash\nsource ~/ros2_curri/agv_ws/install/setup.bash\nros2 pkg executables agv_cpp_examples")],
        "run": "cd ~/ros2_curri/agv_ws\nsource /opt/ros/jazzy/setup.bash\ncolcon build --symlink-install\nsource install/setup.bash\nros2 run agv_cpp_examples status_publisher",
        "validate": "pwd\nfind src -maxdepth 2 -name package.xml -o -name CMakeLists.txt\nros2 pkg prefix agv_cpp_examples",
        "capture": "colcon list --names-only; ros2 pkg executables agv_cpp_examples",
        "errors": [("패키지를 찾을 수 없음", "~/ros2_curri/agv_ws에서 build했는지와 install/setup.bash source를 확인한다."), ("중복 패키지 오류", "src 안에 같은 package name이 두 개 없는지 확인한다.")],
        "mini": "C++ publisher의 message 파라미터를 바꾸고 /cpp_status에서 바뀐 문자열을 확인한다.",
    },
    {
        "id": "M03", "block": "A", "folder": "blocks/A_ros2_basics/M03_pubsub",
        "filename": "M03_Publisher_Subscriber_직접_작성.pptx", "title": "Publisher / Subscriber 직접 작성",
        "goal": "Python timer·publisher·subscriber를 작성해 counter와 Twist를 관찰한다.",
        "completion": "/counter와 /cmd_vel의 값·message type·연결 수를 terminal에서 확인한다.",
        "start": "M02의 agv_control ament_python 패키지가 source 후 인식되는 상태.",
        "previous": "M02 agv_control 패키지", "next": "M04 TF2 frame",
        "flow": ["timer", "publisher\n/counter", "subscriber\ncallback"],
        "files": ["agv_control/agv_control/counter_publisher.py", "agv_control/agv_control/counter_monitor.py", "agv_control/agv_control/cmd_test_node.py", "agv_control/agv_control/velocity_monitor.py", "agv_control/setup.py"],
        "actions": [("console_scripts를 등록한다", "ros2 run은 setup.py에 등록된 실행 이름을 사용한다.", "cd ~/ros2_curri/agv_ws\ncolcon build --symlink-install --packages-select agv_control\nsource install/setup.bash"), ("counter 두 노드를 실행한다", "publisher 출력과 monitor 수신을 분리해 확인한다.", "# 터미널 1\nros2 run agv_control counter_publisher\n# 터미널 2\nros2 run agv_control counter_monitor")],
        "run": "cd ~/ros2_curri/agv_ws && source install/setup.bash\nros2 run agv_control counter_publisher\n# 새 터미널: ros2 run agv_control counter_monitor",
        "validate": "ros2 topic info /counter\nros2 topic echo /counter --once\nros2 interface show geometry_msgs/msg/Twist",
        "capture": "ros2 pkg executables agv_control; ros2 interface show geometry_msgs/msg/Twist",
        "errors": [("No executable found", "setup.py entry_points·재빌드·source 순서로 확인한다."), ("callback이 오지 않음", "토픽 이름과 message type이 양쪽에서 같은지 확인한다.")],
        "mini": "cmd_test_node의 linear.x와 angular.z를 각각 바꾸고 velocity_monitor 출력을 비교한다.",
    },
    {
        "id": "M04", "block": "A", "folder": "blocks/A_ros2_basics/M04_tf2_frames",
        "filename": "M04_TF2와_로봇_좌표계.pptx", "title": "TF2와 로봇 좌표계",
        "goal": "map→odom→base_link→sensor frame의 부모·자식 관계와 Fixed Frame을 이해한다.",
        "completion": "정적 TF를 발행하고 frame 이름과 RViz Fixed Frame의 관계를 설명한다.",
        "start": "M03에서 node/topic을 확인할 수 있는 상태. 아직 실제 AGV 모델은 필요 없다.",
        "previous": "M03 node/topic 관찰", "next": "M05 URDF link/joint",
        "flow": ["map", "odom", "base_link", "lidar_link"],
        "files": [],
        "actions": [("정적 TF를 발행한다", "고정 장착 센서는 시간에 따라 바뀌지 않는 transform을 사용한다.", "source /opt/ros/jazzy/setup.bash\nros2 run tf2_ros static_transform_publisher 0.12 0 0.28 0 0 0 base_link lidar_link"), ("TF 연결을 검사한다", "frame의 부모·자식과 누락 여부를 파일로 확인한다.", "ros2 run tf2_tools view_frames\n# frames.pdf 생성 뒤 base_link → lidar_link 연결을 확인")],
        "run": "ros2 run tf2_ros static_transform_publisher 0.12 0 0.28 0 0 0 base_link lidar_link",
        "validate": "ros2 topic echo /tf_static --once\nros2 run tf2_tools view_frames",
        "capture": "ros2 pkg prefix tf2_ros; ros2 pkg prefix tf2_tools",
        "errors": [("RViz에 아무것도 안 보임", "Fixed Frame과 TF 존재 여부를 먼저 확인한다."), ("센서가 엉뚱한 방향", "rpy 단위·parent/child 순서·축 방향을 점검한다.")],
        "mini": "camera_link의 xyz를 0.30 0 0.26으로 바꾸어 발행하고 lidar_link와 위치 차이를 설명한다.",
    },
    {
        "id": "M05", "block": "B", "folder": "blocks/B_robot_build/M05_urdf_basic_model",
        "filename": "M05_URDF로_기본_AGV_제작.pptx", "title": "URDF로 기본 AGV 제작",
        "goal": "box·cylinder·sphere만으로 body·wheel·caster·sensor frame을 갖춘 AGV를 만든다.",
        "completion": "RobotModel과 TF에서 base, 두 wheel, camera, lidar, imu 연결을 확인한다.",
        "start": "M04 TF2 개념을 이해하고 agv_description 패키지가 존재하는 상태.",
        "previous": "M04 frame tree", "next": "M06 Xacro 모듈화",
        "flow": ["base_link\nbox", "wheel\ncylinder", "sensor\nfixed joint"],
        "files": ["agv_description/curriculum_stages/M05/agv.urdf"],
        "actions": [("단일 URDF를 검사한다", "M05에서는 macro·property·xacro 없이 한 파일에 link와 joint를 직접 작성한다.", "cd ~/ros2_curri/agv_ws\ncheck_urdf src/agv_description/curriculum_stages/M05/agv.urdf"), ("RViz로 표시한다", "RobotModel에는 단일 URDF 텍스트를 robot_description으로 넣고 TF를 함께 확인한다.", "ros2 run robot_state_publisher robot_state_publisher \\\n  --ros-args -p robot_description:=\"$(cat src/agv_description/curriculum_stages/M05/agv.urdf)\"\n# 새 터미널: rviz2 → Add → RobotModel, TF")],
        "run": "cd ~/ros2_curri/agv_ws\ncheck_urdf src/agv_description/curriculum_stages/M05/agv.urdf\n# 다음 M06에서 처음으로 xacro를 사용한다.",
        "validate": "check_urdf ~/ros2_curri/agv_ws/src/agv_description/curriculum_stages/M05/agv.urdf\nrg -n '<xacro:|<link |<joint ' ~/ros2_curri/agv_ws/src/agv_description/curriculum_stages/M05/agv.urdf",
        "capture": "check_urdf src/agv_description/curriculum_stages/M05/agv.urdf",
        "errors": [("XML parsing error", "닫는 태그·따옴표·중첩 구조를 확인한다."), ("바퀴가 눕거나 TF가 없음", "cylinder rpy·joint axis·parent/child link 이름을 확인한다.")],
        "mini": "wheel_radius를 0.08 m에서 0.10 m로 바꾸고 RobotModel에서 바퀴 크기 변화를 확인한다.",
    },
    {
        "id": "M06", "block": "B", "folder": "blocks/B_robot_build/M06_xacro_modularization",
        "filename": "M06_Xacro로_구조_모듈화.pptx", "title": "Xacro로 구조 모듈화",
        "goal": "property·macro·include로 반복되는 wheel·sensor 구조를 한 번만 작성한다.",
        "completion": "wheel_radius와 wheel_base 변경이 좌우 바퀴·충돌 형상에 함께 적용된다.",
        "start": "M05의 단일 AGV 구조가 RViz에 정상 표시되는 상태.",
        "previous": "M05 완성 URDF", "next": "M07 Gazebo 물리",
        "flow": ["agv.urdf.xacro", "include\nmacro", "생성 URDF"],
        "files": ["agv_description/curriculum_stages/M06/agv.urdf.xacro"],
        "actions": [("처음으로 Xacro macro를 만든다", "M05 단일 URDF의 반복 바퀴 구조를 macro/property로 바꾸는 첫 단계다.", "cd ~/ros2_curri/agv_ws\nxacro src/agv_description/curriculum_stages/M06/agv.urdf.xacro > /tmp/agv_from_xacro.urdf"), ("property 변경을 검증한다", "하나의 wheel_radius가 좌우 바퀴에 함께 적용되어야 한다.", "check_urdf /tmp/agv_from_xacro.urdf\n# wheel_radius 값을 바꾼 뒤 같은 두 명령을 다시 실행")],
        "run": "xacro ~/ros2_curri/agv_ws/src/agv_description/curriculum_stages/M06/agv.urdf.xacro > /tmp/agv.urdf\ncheck_urdf /tmp/agv.urdf",
        "validate": "xacro ~/ros2_curri/agv_ws/src/agv_description/curriculum_stages/M06/agv.urdf.xacro > /tmp/agv.urdf\nrg -n 'wheel_radius|left_wheel_joint|right_wheel_joint' /tmp/agv.urdf",
        "capture": "xacro src/agv_description/curriculum_stages/M06/agv.urdf.xacro > /tmp/m06.urdf && check_urdf /tmp/m06.urdf",
        "errors": [("unknown macro/property", "선언 순서·xacro namespace·철자를 확인한다."), ("include file not found", "패키지 이름과 설치 대상·파일 경로를 확인한다.")],
        "mini": "wheel_radius 0.08→0.12 m를 적용한 뒤 joint origin과 cylinder radius가 함께 바뀌는지 비교한다.",
    },
    {
        "id": "M07", "block": "B", "folder": "blocks/B_robot_build/M07_physics_modeling",
        "filename": "M07_Gazebo용_물리_요소_추가.pptx", "title": "Gazebo용 물리 요소 추가",
        "goal": "visual·collision·inertial·friction이 각각 물리 시뮬레이션에 미치는 영향을 구분한다.",
        "completion": "AGV가 바닥에 안정적으로 놓이고 바퀴 접촉·충돌 파라미터가 모델에 존재한다.",
        "start": "M06 Xacro가 RViz에서 보이지만 Gazebo 물리 모델은 아직 별도인 상태.",
        "previous": "M06 Xacro 치수", "next": "M08 World/Spawn",
        "flow": ["visual\n보이는 모양", "collision\n접촉", "inertial\n질량/관성", "friction\n접지"],
        "files": ["agv_description/curriculum_stages/M07/model.sdf"],
        "actions": [("SDF 물리 모델을 검사한다", "M07에는 visual·collision·inertial·friction만 있으며 sensor와 drive plugin은 아직 없다.", "cd ~/ros2_curri/agv_ws\ngz sdf -k src/agv_description/curriculum_stages/M07/model.sdf"), ("값을 한 번에 하나씩 바꾼다", "증상 원인을 찾으려면 mass→inertia→collision→friction 순서를 지킨다.", "rg -n '<mass>|<inertia>|<collision>|<friction>' src/agv_description/curriculum_stages/M07/model.sdf")],
        "run": "cd ~/ros2_curri/agv_ws\ngz sdf -k src/agv_description/curriculum_stages/M07/model.sdf",
        "validate": "gz sdf -k ~/ros2_curri/agv_ws/src/agv_description/curriculum_stages/M07/model.sdf\nrg -n '<sensor|DiffDrive' ~/ros2_curri/agv_ws/src/agv_description/curriculum_stages/M07/model.sdf",
        "capture": "gz sdf -k src/agv_description/curriculum_stages/M07/model.sdf; rg -n '<mass>|<friction>|left_wheel_joint' src/agv_description/curriculum_stages/M07/model.sdf",
        "errors": [("spawn 직후 떨리거나 날아감", "inertia·collision overlap·너무 작은 mass를 확인한다."), ("바퀴가 헛돎", "friction·contact·joint axis·wheel radius를 확인한다.")],
        "mini": "wheel friction mu를 1.0에서 0.3으로만 바꾼 뒤 직진 시 미끄럼 변화를 기록한다.",
    },
    {
        "id": "M08", "block": "B", "folder": "blocks/B_robot_build/M08_gazebo_world_spawn",
        "filename": "M08_Gazebo_World와_Robot_Spawn.pptx", "title": "Gazebo World와 Robot Spawn",
        "goal": "warehouse World와 AGV model을 package 경로로 열고 충돌 환경을 준비한다.",
        "completion": "ground·벽·박스가 있는 World에 AGV가 spawn되어 있는 Gazebo 화면을 확인한다.",
        "start": "M07의 SDF AGV 물리 모델이 유효한 상태.",
        "previous": "M07 물리 모델", "next": "M09 Differential Drive",
        "flow": ["warehouse.sdf", "ros_gz_sim", "model URI\nAGV spawn"],
        "files": ["agv_description/curriculum_stages/M08/model.config", "agv_description/curriculum_stages/M08/model.sdf", "agv_description/curriculum_stages/M08/warehouse.sdf"],
        "actions": [("model과 World를 각각 검사한다", "M08은 M07 물리 AGV를 World에 배치하는 단계이며 DiffDrive·Camera·LiDAR·IMU는 아직 넣지 않는다.", "cd ~/ros2_curri/agv_ws\ngz sdf -k src/agv_description/curriculum_stages/M08/model.sdf\nrg -n 'model://agv|<world name=' src/agv_description/curriculum_stages/M08/warehouse.sdf"), ("World를 실행한다", "M08 snapshot을 agv_gazebo models/worlds 경로에 복사한 뒤 Gazebo를 실행해 spawn만 확인한다.", "gz sim -r src/agv_description/curriculum_stages/M08/warehouse.sdf\n# 이 단계에는 ROS bridge와 구동 topic이 없다.")],
        "run": "cd ~/ros2_curri/agv_ws\ngz sim -r src/agv_description/curriculum_stages/M08/warehouse.sdf",
        "validate": "gz sdf -k ~/ros2_curri/agv_ws/src/agv_description/curriculum_stages/M08/model.sdf\nrg -n '<sensor|DiffDrive|/cmd_vel|/odom' ~/ros2_curri/agv_ws/src/agv_description/curriculum_stages/M08/model.sdf",
        "capture": "gz sdf -k src/agv_description/curriculum_stages/M08/model.sdf; rg -n 'model://agv|<world name=' src/agv_description/curriculum_stages/M08/warehouse.sdf",
        "visual": "blocks/B_robot_build/captures/01_gazebo_world_spawn_actual.png", "visual_caption": "실제 Gazebo Sim에서 warehouse World에 spawn된 AGV",
        "errors": [("world 파일을 못 찾음", "launch의 package share path와 install data_files를 확인한다."), ("모델이 공중/바닥 아래", "초기 z·collision origin·geometry 크기를 확인한다.")],
        "mini": "World의 box pose x 좌표를 1 m 바꾼 후 저장·재실행하여 top-view 위치 변화를 확인한다.",
    },
    {
        "id": "M09", "block": "C", "folder": "blocks/C_drive_visualization/M09_differential_drive",
        "filename": "M09_Differential_Drive_구동.pptx", "title": "Differential Drive 구동",
        "goal": "좌우 바퀴 속도 차이와 Twist linear.x/angular.z로 직진·회전·곡선 주행을 설명한다.",
        "completion": "/cmd_vel에서 odom과 바퀴 구동 topic으로 이어지는 경로를 확인한다.",
        "start": "M08에서 AGV가 World에 안정적으로 spawn된 상태.",
        "previous": "M08 World/AGV", "next": "M10 ros_gz_bridge",
        "flow": ["/cmd_vel\nTwist", "DiffDrive", "left/right\nwheel", "/odom"],
        "files": ["agv_gazebo/models/agv/model.sdf", "agv_control/agv_control/cmd_test_node.py"],
        "actions": [("직진 명령을 보낸다", "linear.x는 m/s이고 angular.z는 rad/s다.", "ros2 topic pub --rate 10 /cmd_vel geometry_msgs/msg/Twist \\\n  \"{linear: {x: 0.15}, angular: {z: 0.0}}\""), ("회전과 odom을 확인한다", "제자리 회전은 linear.x=0, angular.z만 양수/음수로 바꾼다.", "ros2 topic pub --rate 10 /cmd_vel geometry_msgs/msg/Twist \\\n  \"{linear: {x: 0.0}, angular: {z: 0.4}}\"\nros2 topic echo /odom --once")],
        "run": "ros2 launch agv_gazebo gazebo.launch.py\n# 새 터미널: ros2 run agv_control cmd_test_node",
        "validate": "ros2 topic info /cmd_vel\nros2 topic echo /odom --once",
        "capture": "rg -n 'DiffDrive|wheel_separation|wheel_radius|<topic>/cmd_vel' src/agv_gazebo/models/agv/model.sdf",
        "visual": "blocks/C_drive_visualization/captures/01_gazebo_drive_actual.png", "visual_caption": "실제 Gazebo Sim: World 안의 AGV와 장애물·목표물 배치",
        "errors": [("바퀴는 도는데 이동 안 함", "joint axis·friction·contact·wheel radius를 확인한다."), ("회전 방향이 반대", "left/right joint와 angular.z 부호를 확인한다.")],
        "mini": "0.15 m/s 직진, 0.4 rad/s 회전, 둘 다 비영 값인 곡선을 각각 2초씩 시험해 경로를 비교한다.",
    },
    {
        "id": "M10", "block": "C", "folder": "blocks/C_drive_visualization/M10_ros_gz_bridge",
        "filename": "M10_ros_gz_bridge_ROS2_Gazebo_연결.pptx", "title": "ros_gz_bridge로 ROS 2 ↔ Gazebo 연결",
        "goal": "ROS 2 message와 Gazebo Transport message 사이의 방향·타입·이름 매핑을 만든다.",
        "completion": "/cmd_vel ROS→Gazebo 및 odom/scan/imu/camera/clock Gazebo→ROS 매핑을 확인한다.",
        "start": "Gazebo 내부에 AGV와 sensor topic이 있고 ROS 2 외부와 구분되는 상태.",
        "previous": "M09 /cmd_vel", "next": "M11 RViz2",
        "flow": ["ROS 2\n/cmd_vel", "ros_gz_bridge", "Gazebo\nDiffDrive"],
        "files": ["agv_gazebo/config/bridge.yaml", "agv_gazebo/launch/gazebo.launch.py"],
        "actions": [("bridge YAML을 작성한다", "방향과 양쪽 message type을 한 topic마다 정확히 적는다.", "cd ~/ros2_curri/agv_ws\nnano src/agv_gazebo/config/bridge.yaml"), ("launch가 YAML을 읽게 한다", "config_file은 package share 경로에서 찾도록 launch에 넣는다.", "colcon build --symlink-install --packages-select agv_gazebo\nsource install/setup.bash\nros2 launch agv_gazebo gazebo.launch.py")],
        "run": "source ~/ros2_curri/agv_ws/install/setup.bash\nros2 launch agv_gazebo gazebo.launch.py",
        "validate": "ros2 topic list | rg '/cmd_vel|/odom|/scan|/imu|/camera|/clock'\nros2 topic info /scan",
        "capture": "ros2 pkg prefix ros_gz_bridge; printf '\\n--- bridge.yaml ---\\n'; cat src/agv_gazebo/config/bridge.yaml",
        "errors": [("ROS topic이 생성되지 않음", "bridge 프로세스·topic 이름·message type을 확인한다."), ("한 방향으로만 전달됨", "direction ROS_TO_GZ/GZ_TO_ROS 설정을 확인한다.")],
        "mini": "bridge.yaml에서 /scan 항목의 ros_type_name을 읽고 `ros2 interface show sensor_msgs/msg/LaserScan`으로 필드를 확인한다.",
    },
    {
        "id": "M11", "block": "C", "folder": "blocks/C_drive_visualization/M11_rviz_integration",
        "filename": "M11_RViz2_구동_상태_검증.pptx", "title": "RViz2에서 구동 상태 검증",
        "goal": "Gazebo와 RViz2 역할을 분리하고 RobotModel·TF·센서 display를 같은 frame에서 검증한다.",
        "completion": "agv.rviz에 RobotModel·TF·LaserScan·Odometry 설정을 저장하고 실제 화면을 확인한다.",
        "start": "M10 bridge 뒤에 robot_description·TF·odom·scan이 ROS 2에 존재하는 상태.",
        "previous": "M10 bridge topic", "next": "M12~M15 센서 display",
        "flow": ["robot_description", "TF / odom", "RViz Displays"],
        "files": [],
        "actions": [("RViz를 설정 파일로 연다", "같은 display 설정을 다음 센서 모듈에서도 재사용한다.", "source ~/ros2_curri/agv_ws/install/setup.bash\nrviz2 -d ~/ros2_curri/agv_ws/src/agv_description/rviz/agv.rviz"), ("RobotModel과 TF를 먼저 확인한다", "모델과 frame이 맞아야 scan·image 위치도 맞는다.", "Fixed Frame: odom 또는 base_link\nAdd → RobotModel → Description Topic: /robot_description\nAdd → TF"), ("센서 display를 추가한다", "각 display는 topic과 message frame_id를 TF로 변환한다.", "Add → LaserScan → Topic: /scan\nAdd → Odometry → Topic: /odom\nFile → Save Config As → agv.rviz")],
        "run": "source ~/ros2_curri/agv_ws/install/setup.bash\nros2 launch agv_bringup agv_sim.launch.py",
        "validate": "ros2 topic info /robot_description\nros2 topic echo /scan --once\nros2 topic echo /odom --once",
        "capture": "ros2 pkg prefix rviz2; rg -n 'Fixed Frame|RobotModel|LaserScan|Odometry' src/agv_description/rviz/agv.rviz",
        "visual": "blocks/C_drive_visualization/captures/01_rviz_integrated_actual.png", "visual_caption": "실제 RViz2: RobotModel·TF·Odometry·LaserScan·Path·Marker display를 함께 연 상태",
        "errors": [("RobotModel이 안 보임", "robot_description·Fixed Frame·TF를 순서대로 확인한다."), ("LaserScan 위치가 이상함", "scan header.frame_id와 lidar TF를 확인한다.")],
        "mini": "Fixed Frame을 의도적으로 없는 이름으로 바꾼 뒤 Global Status 오류를 읽고 odom으로 되돌린다.",
    },
    {
        "id": "M12", "block": "D", "folder": "blocks/D_sensors/M12_camera_sensor",
        "filename": "M12_Camera_Sensor.pptx", "title": "Camera Sensor",
        "goal": "Gazebo camera의 FOV·해상도·update_rate가 ROS Image topic으로 연결되는 흐름을 만든다.",
        "completion": "/camera/image_raw가 bridge를 통과하고 cv_bridge 소비 노드가 받을 준비가 된 상태를 확인한다.",
        "start": "M11의 camera_link와 RViz display 환경이 준비된 상태.",
        "previous": "M11 RViz2", "next": "M13 LiDAR",
        "flow": ["Gazebo camera", "bridge", "/camera/image_raw", "cv_bridge"],
        "files": ["agv_gazebo/models/agv/model.sdf", "agv_gazebo/config/bridge.yaml", "agv_vision/agv_vision/yolo_node.py"],
        "actions": [("camera sensor 설정을 확인한다", "해상도·FOV·update_rate는 생성되는 영상의 크기와 주기를 정한다.", "rg -n '<sensor name=\"camera\"|<width>|<height>|<horizontal_fov>|<update_rate>' \\\n  ~/ros2_curri/agv_ws/src/agv_gazebo/models/agv/model.sdf"), ("Image topic을 확인한다", "topic 이름은 SDF sensor와 bridge YAML 양쪽에서 일치해야 한다.", "ros2 topic info /camera/image_raw\nros2 topic hz /camera/image_raw\nros2 run rqt_image_view rqt_image_view")],
        "run": "ros2 launch agv_bringup agv_sim.launch.py rviz:=false autonomy:=false\n# 새 터미널: ros2 run rqt_image_view rqt_image_view /camera/image_raw",
        "validate": "ros2 topic info /camera/image_raw\nros2 topic hz /camera/image_raw",
        "capture": "rg -n '<sensor name=\"camera\"|<topic>/camera/image_raw|<width>|<height>' src/agv_gazebo/models/agv/model.sdf; rg -n 'camera/image_raw' src/agv_gazebo/config/bridge.yaml",
        "visual": "blocks/D_sensors/captures/02_camera_image_actual.png", "visual_caption": "실제 /camera/image_raw 캡처: 화면 방향·FOV·world 배치를 점검하는 기준",
        "errors": [("영상이 검정", "camera pose/FOV·world light·clipping range를 확인한다."), ("RViz 이미지가 비어 있음", "Image topic·QoS·bridge mapping을 확인한다.")],
        "mini": "camera update_rate 15→5 Hz와 width 640→320을 각각 바꾸고 topic hz와 처리량 차이를 기록한다.",
    },
    {
        "id": "M13", "block": "D", "folder": "blocks/D_sensors/M13_lidar_sensor",
        "filename": "M13_2D_LiDAR_LaserScan.pptx", "title": "2D LiDAR / LaserScan",
        "goal": "LaserScan 범위 배열에서 전방 ±15° 유효 거리의 최소값을 계산한다.",
        "completion": "/scan frame·RViz point·/obstacle_distance 출력이 장애물 위치에 반응한다.",
        "start": "lidar_link가 있고 World에 벽·박스가 배치된 상태.",
        "previous": "M12 Camera pipeline", "next": "M14 IMU",
        "flow": ["gpu_lidar", "bridge\n/scan", "lidar_processor", "/obstacle_distance"],
        "files": ["agv_gazebo/models/agv/model.sdf", "agv_gazebo/config/bridge.yaml", "agv_sensors/agv_sensors/lidar_processor.py"],
        "actions": [("LiDAR 범위와 angle을 읽는다", "ranges의 index는 angle_min + index×angle_increment로 실제 각도가 된다.", "ros2 interface show sensor_msgs/msg/LaserScan\nros2 topic echo /scan --once"), ("전방 최소거리 노드를 실행한다", "inf·NaN은 충돌 판단에서 제외해야 안전한 최소값을 얻는다.", "source ~/ros2_curri/agv_ws/install/setup.bash\nros2 run agv_sensors lidar_processor\nros2 topic echo /obstacle_distance")],
        "run": "ros2 launch agv_bringup agv_sim.launch.py rviz:=false\n# 새 터미널: ros2 run agv_sensors lidar_processor",
        "validate": "ros2 topic info /scan\nros2 topic hz /scan\nros2 topic echo /obstacle_distance --once",
        "capture": "ros2 pkg executables agv_sensors; rg -n '<sensor name=\"lidar\"|<samples>|<min_angle>|<max_angle>' src/agv_gazebo/models/agv/model.sdf",
        "visual": "blocks/D_sensors/captures/03_rviz_sensors_actual.png", "visual_caption": "실제 RViz2 통합 display: LaserScan·frame·Fixed Frame을 점검하는 화면",
        "errors": [("/scan이 없음", "SDF sensor → bridge YAML → ROS topic 순서로 확인한다."), ("최소거리가 항상 inf", "유효 range 필터와 sector index 계산을 확인한다.")],
        "mini": "front_half_angle_deg를 15°와 45°로 각각 실행해 장애물 판정 영역이 어떻게 달라지는지 비교한다.",
    },
    {
        "id": "M14", "block": "D", "folder": "blocks/D_sensors/M14_imu_sensor",
        "filename": "M14_IMU_Sensor.pptx", "title": "IMU Sensor",
        "goal": "angular_velocity·linear_acceleration·frame_id를 AGV 회전과 가속의 좌표축으로 해석한다.",
        "completion": "/imu/data의 z축 gyro 값과 noise 설정을 terminal에서 확인한다.",
        "start": "IMU sensor가 base_link에 장착되어 있고 M09 주행이 가능한 상태.",
        "previous": "M13 /scan", "next": "M15 운영/QoS",
        "flow": ["Gazebo IMU", "bridge", "/imu/data", "imu_monitor"],
        "files": ["agv_gazebo/models/agv/model.sdf", "agv_gazebo/worlds/warehouse.sdf", "agv_gazebo/config/bridge.yaml", "agv_sensors/agv_sensors/imu_monitor.py"],
        "actions": [("Imu message 필드를 확인한다", "z축 angular velocity는 제자리 회전에서 가장 먼저 변화를 보인다.", "ros2 interface show sensor_msgs/msg/Imu\nros2 topic echo /imu/data --once"), ("IMU monitor를 실행한다", "출력에 단위와 frame을 함께 남기면 축 오류를 찾기 쉽다.", "source ~/ros2_curri/agv_ws/install/setup.bash\nros2 run agv_sensors imu_monitor")],
        "run": "ros2 launch agv_bringup agv_sim.launch.py rviz:=false\n# 새 터미널: ros2 run agv_sensors imu_monitor",
        "validate": "ros2 topic info /imu/data\nros2 topic hz /imu/data\nros2 topic echo /imu/data --once",
        "capture": "ros2 pkg executables agv_sensors; rg -n '<sensor name=\"imu\"|<stddev>|<update_rate>' src/agv_gazebo/models/agv/model.sdf",
        "visual": "blocks/B_robot_build/captures/01_gazebo_world_spawn_actual.png", "visual_caption": "실제 Gazebo AGV: IMU가 base_link 위에 고정 장착되는 물리 모델 화면",
        "errors": [("값 축이 예상과 다름", "imu_link 축·origin rpy·header.frame_id를 확인한다."), ("값이 지나치게 요동", "noise/bias·update_rate·시뮬레이션 부하를 확인한다.")],
        "mini": "정지·직진·제자리 회전에서 angular_velocity.z를 각각 기록하고 부호가 어느 동작에서 바뀌는지 설명한다.",
    },
    {
        "id": "M15", "block": "D", "folder": "blocks/D_sensors/M15_sensor_operations",
        "filename": "M15_센서_운영과_동기화.pptx", "title": "센서 운영과 동기화",
        "goal": "use_sim_time·/clock·QoS·rate·frame_id로 camera/LiDAR/IMU를 공통 점검한다.",
        "completion": "모든 sensor 노드가 같은 시뮬레이션 시간과 기대 publish rate를 사용한다.",
        "start": "M12~M14에서 세 sensor topic이 각각 존재하는 상태.",
        "previous": "M12~M14 센서 topic", "next": "M16 제어 interface",
        "flow": ["/clock", "use_sim_time", "sensor topic", "processing node"],
        "files": ["agv_bringup/config/sensors.yaml", "agv_gazebo/config/bridge.yaml", "agv_sensors/agv_sensors/lidar_processor.py", "agv_sensors/agv_sensors/imu_monitor.py"],
        "actions": [("시뮬레이션 시간을 설정한다", "Gazebo /clock을 안 쓰면 sensor timestamp와 node 시간이 어긋날 수 있다.", "ros2 param get /lidar_processor use_sim_time\nros2 topic echo /clock --once"), ("rate와 QoS를 점검한다", "topic이 있어도 QoS가 호환되지 않으면 callback이 받지 못할 수 있다.", "ros2 topic hz /scan\nros2 topic hz /imu/data\nros2 topic info --verbose /scan")],
        "run": "ros2 launch agv_bringup agv_sim.launch.py rviz:=false",
        "validate": "ros2 topic echo /clock --once\nros2 topic hz /scan\nros2 topic hz /imu/data\nros2 param get /lidar_processor use_sim_time",
        "capture": "cat src/agv_bringup/config/sensors.yaml; printf '\\n--- QoS source ---\\n'; rg -n 'qos_profile_sensor_data|use_sim_time' src/agv_sensors",
        "visual": "blocks/D_sensors/captures/03_rviz_sensors_actual.png", "visual_caption": "실제 RViz2: 여러 sensor display를 같은 Fixed Frame에서 여는 운영 점검 화면",
        "errors": [("ROS 노드 시간이 안 감", "use_sim_time과 /clock bridge를 확인한다."), ("topic은 있는데 subscriber가 못 받음", "publisher/subscriber QoS reliability 호환을 확인한다.")],
        "mini": "/scan과 /imu/data의 hz를 표로 기록하고 SDF update_rate 10/100 Hz와 비교한다.",
    },
    {
        "id": "M16", "block": "E", "folder": "blocks/E_autonomy_logic/M16_ros2_control",
        "filename": "M16_ros2_control_gz_ros2_control.pptx", "title": "ros2_control + gz_ros2_control",
        "goal": "gz_ros2_control이 만든 controller_manager에서 joint_state_broadcaster와 diff_drive_controller를 실제 active로 만든다.",
        "completion": "`ros2 control list_controllers`에서 joint_state_broadcaster와 diff_drive_controller가 모두 active이고 wheel joint state가 publish된다.",
        "start": "M09 native Gazebo DiffDrive가 동작하고 wheel joint 이름을 알고 있는 상태.",
        "previous": "M09 wheel joint", "next": "M17~M20 autonomy command",
        "flow": ["/cmd_vel", "controller_manager", "gz_ros2_control", "Gazebo joints"],
        "files": ["agv_control/config/controllers.yaml", "agv_description/urdf/agv_ros2_control.urdf.xacro", "agv_gazebo/models/agv_ros2_control/model.sdf", "agv_gazebo/launch/ros2_control.launch.py"],
        "actions": [("controller YAML을 작성한다", "controller type과 wheel parameter는 URDF/SDF joint 이름과 같아야 한다.", "nano ~/ros2_curri/agv_ws/src/agv_control/config/controllers.yaml\nros2 pkg prefix gz_ros2_control"), ("실제 controller를 활성화한다", "M16은 M09 native DiffDrive와 별도 launch를 사용한다. plugin이 controller_manager를 만들고 spawner가 두 controller를 active로 전환한다.", "ros2 launch agv_gazebo ros2_control.launch.py\n# 새 터미널\nros2 control list_controllers -c /controller_manager")],
        "run": "source ~/ros2_curri/agv_ws/install/setup.bash\nros2 launch agv_gazebo ros2_control.launch.py\n# 새 터미널: ros2 control list_controllers -c /controller_manager",
        "validate": "ros2 control list_controllers -c /controller_manager\nros2 topic echo /joint_states --once",
        "capture": "ros2 pkg prefix gz_ros2_control; ros2 pkg prefix controller_manager; cat src/agv_control/config/controllers.yaml",
        "visual": "blocks/E_autonomy_logic/captures/01_controller_active_actual.png", "visual_caption": "실제 M16 격리 launch의 ros2 control 출력: 두 controller가 모두 active",
        "errors": [("controller가 inactive", "load/activate 순서와 YAML controller 이름을 확인한다."), ("joint를 찾을 수 없음", "URDF/SDF joint 이름과 controller 목록을 한 글자씩 비교한다.")],
        "mini": "controllers.yaml의 wheel_radius와 model.sdf의 wheel_radius가 모두 0.08 m인지 확인하고 한쪽만 바꿨을 때 생길 오차를 설명한다.",
    },
    {
        "id": "M17", "block": "E", "folder": "blocks/E_autonomy_logic/M17_yolo_vision",
        "filename": "M17_Camera_YOLO_Vision_Node.pptx", "title": "Camera + YOLO Vision Node",
        "goal": "ROS Image→cv_bridge→YOLO→DetectionArray 흐름과 모델·threshold 재현 정보를 관리한다.",
        "completion": "/detections message 형식과 YOLO enable_yolo/model_path/threshold parameter를 확인한다.",
        "start": "M12에서 /camera/image_raw가 안정적으로 publish되고 cv_bridge 의존성이 설치된 상태.",
        "previous": "M12 camera image", "next": "M18 safety fusion",
        "flow": ["Image", "cv_bridge", "YOLO", "/detections"],
        "files": ["agv_interfaces/msg/Detection.msg", "agv_interfaces/msg/DetectionArray.msg", "agv_interfaces/CMakeLists.txt", "agv_vision/agv_vision/yolo_node.py", "agv_bringup/config/vision.yaml", "agv_vision/setup.py"],
        "actions": [("interface를 먼저 빌드한다", "vision node보다 custom message package가 먼저 install/source되어야 import된다.", "cd ~/ros2_curri/agv_ws\ncolcon build --symlink-install --packages-select agv_interfaces agv_vision\nsource install/setup.bash\nros2 interface show agv_interfaces/msg/Detection"), ("YOLO 실행 정책을 정한다", "모델 파일과 Python dependency는 교육 환경에서 고정해야 재현된다.", "ros2 run agv_vision yolo_node --ros-args \\\n  -p enable_yolo:=false\n# ultralytics와 모델 파일 준비 뒤에만 true로 변경")],
        "run": "source ~/ros2_curri/agv_ws/install/setup.bash\nros2 run agv_vision yolo_node --ros-args -p enable_yolo:=false",
        "validate": "ros2 interface show agv_interfaces/msg/Detection\nros2 topic info /detections\nros2 param get /yolo_node confidence_threshold",
        "capture": "ros2 interface show agv_interfaces/msg/Detection; ros2 pkg executables agv_vision; cat src/agv_bringup/config/vision.yaml",
        "visual": "blocks/E_autonomy_logic/captures/01_vision_debug_actual.png", "visual_caption": "실제 /vision/debug_image: fallback이 ‘target not visible’을 표시해 camera pose·FOV·world 배치를 점검하는 오류 사례",
        "errors": [("모델 파일을 못 찾음", "model_path와 모델 파일 배포 정책을 확인한다."), ("Custom msg import 실패", "agv_interfaces 빌드·source·package dependency를 확인한다.")],
        "mini": "confidence_threshold를 0.50에서 0.70으로 바꿨을 때 어떤 검출이 제거될지 설명하고 parameter 값을 확인한다.",
    },
    {
        "id": "M18", "block": "E", "folder": "blocks/E_autonomy_logic/M18_obstacle_perception",
        "filename": "M18_LiDAR_장애물_감지와_Perception_결합.pptx", "title": "LiDAR 장애물 감지와 Perception 결합",
        "goal": "전방 LiDAR 거리와 mission command 사이에 Safety Controller를 두어 STOP을 최우선으로 만든다.",
        "completion": "0.5 m 안 장애물에서 forward cmd_vel이 zero Twist로 바뀌는 조건을 설명한다.",
        "start": "M13 obstacle_distance와 M17 DetectionArray가 각각 topic으로 나오는 상태.",
        "previous": "M13 + M17 outputs", "next": "M19 FSM",
        "flow": ["/scan", "safety\npriority", "/cmd_vel", "DiffDrive"],
        "files": ["agv_sensors/agv_sensors/lidar_processor.py", "agv_control/agv_control/safety_controller.py", "agv_bringup/config/mission.yaml"],
        "actions": [("안전 임계값을 고정한다", "distance는 m 단위이며 safety가 vision보다 먼저 final /cmd_vel을 결정한다.", "ros2 run agv_control safety_controller --ros-args \\\n  -p stop_distance:=0.5 -p front_half_angle_deg:=15.0"), ("raw와 safe command를 분리한다", "mission은 /cmd_vel_raw, 안전 필터만 /cmd_vel을 발행한다.", "ros2 topic echo /cmd_vel_raw\nros2 topic echo /cmd_vel\nros2 topic echo /scan --once")],
        "run": "source ~/ros2_curri/agv_ws/install/setup.bash\nros2 run agv_control safety_controller --ros-args -p stop_distance:=0.5",
        "validate": "ros2 topic info /cmd_vel_raw\nros2 topic info /cmd_vel\nros2 param get /safety_controller stop_distance",
        "capture": "ros2 pkg executables agv_control; rg -n 'stop_distance|cmd_vel_raw|/cmd_vel|front_half_angle' src/agv_control/agv_control/safety_controller.py",
        "visual": "blocks/E_autonomy_logic/captures/02_rviz_marker_path_actual.png", "visual_caption": "실제 RViz2: safety sector Marker와 Path를 추가해 확인하는 통합 display",
        "errors": [("장애물이 있는데 멈추지 않음", "sector 각도·finite range·threshold m 단위를 확인한다."), ("항상 STOP", "inf/NaN 처리와 lidar frame 전방 방향을 확인한다.")],
        "mini": "stop_distance 0.50→0.80 m로 바꾸면 어떤 상황에서 더 일찍 멈추는지 /cmd_vel 로그 기준으로 작성한다.",
    },
    {
        "id": "M19", "block": "E", "folder": "blocks/E_autonomy_logic/M19_mission_fsm",
        "filename": "M19_Mission_State_Machine.pptx", "title": "Mission State Machine",
        "goal": "SEARCH·APPROACH·AVOID·GOAL 상태를 입력 조건과 속도 명령으로 명시한다.",
        "completion": "/mission_state 로그와 /cmd_vel_raw가 detection·obstacle input에 따라 바뀌는 구조를 확인한다.",
        "start": "M18 safety와 M17 target detection input이 준비된 상태.",
        "previous": "M18 priority decision", "next": "M20 PID",
        "flow": ["SEARCH", "APPROACH", "AVOID", "GOAL"],
        "files": ["agv_mission/agv_mission/mission_manager.py", "agv_bringup/config/mission.yaml", "agv_mission/setup.py"],
        "actions": [("FSM node를 빌드하고 실행한다", "상태는 흩어진 if문이 아니라 한 곳에서 /cmd_vel_raw를 결정한다.", "cd ~/ros2_curri/agv_ws\ncolcon build --symlink-install --packages-select agv_mission\nsource install/setup.bash\nros2 run agv_mission mission_manager"), ("상태 변화를 관찰한다", "상태의 이름과 속도 command를 서로 다른 topic에서 확인한다.", "ros2 topic echo /mission_state\nros2 topic echo /cmd_vel_raw")],
        "run": "source ~/ros2_curri/agv_ws/install/setup.bash\nros2 run agv_mission mission_manager",
        "validate": "ros2 topic info /mission_state\nros2 topic info /cmd_vel_raw\nros2 param get /mission_manager stop_distance",
        "capture": "ros2 pkg executables agv_mission; cat src/agv_bringup/config/mission.yaml",
        "visual": "blocks/E_autonomy_logic/captures/02_rviz_marker_path_actual.png", "visual_caption": "실제 RViz2: mission text Marker를 같은 scene에서 확인하는 display 구성",
        "errors": [("상태가 전환되지 않음", "callback input·조건식·거리 단위를 확인한다."), ("SEARCH↔APPROACH 반복", "detection timeout·hysteresis·우선순위를 추가 점검한다.")],
        "mini": "target 없음, target 있음, obstacle<0.5 m의 세 입력에 대해 예상 state와 /cmd_vel_raw를 표로 작성한다.",
    },
    {
        "id": "M20", "block": "E", "folder": "blocks/E_autonomy_logic/M20_pid_driving_control",
        "filename": "M20_PID와_주행_제어.pptx", "title": "PID와 주행 제어",
        "goal": "target error를 안전한 linear/angular Twist로 제한하는 제어 시작점을 만든다.",
        "completion": "kp_angular·max speed parameter와 saturation이 /cmd_vel_raw 출력을 제한하는 것을 확인한다.",
        "start": "M19가 상태별 desired command를 만들고 target 중심 오차를 계산할 수 있는 상태.",
        "previous": "M19 state command", "next": "M21 YAML/Launch",
        "flow": ["target error", "P / PID\ncontroller", "/cmd_vel_raw", "Safety"],
        "files": ["agv_control/agv_control/pid_controller.py", "agv_control/setup.py", "agv_bringup/config/mission.yaml"],
        "actions": [("제어 파라미터를 확인한다", "P gain과 최대 속도는 overshoot 전에 먼저 보수적으로 제한한다.", "ros2 run agv_control pid_controller --ros-args \\\n  -p kp_angular:=0.004 -p max_linear_speed:=0.25"), ("error와 command를 분리해 관찰한다", "입력 부호가 바뀌면 angular.z 부호도 기대 방향으로 바뀌어야 한다.", "ros2 topic echo /target_error\nros2 topic echo /cmd_vel_raw")],
        "run": "source ~/ros2_curri/agv_ws/install/setup.bash\nros2 run agv_control pid_controller --ros-args -p kp_angular:=0.004",
        "validate": "ros2 topic info /target_error\nros2 topic info /cmd_vel_raw\nros2 param get /pid_controller kp_angular",
        "capture": "ros2 pkg executables agv_control; rg -n 'kp_angular|max_linear_speed|max_angular_speed|cmd_vel_raw' src/agv_control/agv_control/pid_controller.py",
        "visual": "blocks/E_autonomy_logic/captures/02_rviz_marker_path_actual.png", "visual_caption": "실제 RViz2: /path와 robot pose를 함께 보며 PID 전후 궤적을 비교하는 구성",
        "errors": [("로봇이 반대 방향 회전", "error 정의와 angular.z 부호를 확인한다."), ("진동이 심함", "kp/kd·update rate·latency·saturation을 확인한다.")],
        "mini": "kp_angular을 0.002, 0.004, 0.008로 바꿨을 때 반응 속도와 진동 위험을 예측해 기록한다.",
    },
    {
        "id": "M21", "block": "F", "folder": "blocks/F_integration/M21_launch_yaml_rosbag",
        "filename": "M21_Parameter_YAML_Launch_rosbag2.pptx", "title": "Parameter, YAML, Launch, rosbag2",
        "goal": "흩어진 숫자와 수동 실행을 YAML·Launch·rosbag2로 재현 가능한 실행으로 묶는다.",
        "completion": "한 launch 명령·YAML parameter·rosbag record/replay 명령을 설명하고 확인한다.",
        "start": "M20까지 각 노드를 수동 터미널에서 실행할 수 있는 상태.",
        "previous": "M20 controller", "next": "M22 final mission",
        "flow": ["YAML config", "agv_sim.launch.py", "nodes + Gazebo", "rosbag2"],
        "files": ["agv_bringup/config/robot.yaml", "agv_bringup/config/sensors.yaml", "agv_bringup/config/vision.yaml", "agv_bringup/config/mission.yaml", "agv_control/config/controllers.yaml", "agv_bringup/launch/agv_sim.launch.py", "agv_gazebo/launch/gazebo.launch.py"],
        "actions": [("YAML 책임을 분리한다", "robot·sensor·vision·mission·controller 값은 실행 코드와 분리한다.", "find ~/ros2_curri/agv_ws/src -path '*/config/*.yaml' -type f | sort\nros2 param get /mission_manager stop_distance"), ("한 launch로 실행한다", "launch는 Gazebo, TF, processing, mission, safety의 시작 순서를 한 곳에서 관리한다.", "source ~/ros2_curri/agv_ws/install/setup.bash\nros2 launch agv_bringup agv_sim.launch.py"), ("주행 데이터를 기록한다", "bag은 simulator 재실행 없이 processing 결과를 비교할 입력 데이터다.", "ros2 bag record -o ~/agv_bag_01 \\\n  /scan /imu/data /odom /camera/image_raw\n# Ctrl-C 뒤: ros2 bag info ~/agv_bag_01")],
        "run": "cd ~/ros2_curri/agv_ws && source install/setup.bash\nros2 launch agv_bringup agv_sim.launch.py\n# 별도 터미널: ros2 bag record -o ~/agv_bag_01 /scan /imu/data /odom /camera/image_raw",
        "validate": "ros2 launch agv_bringup agv_sim.launch.py --show-args\nros2 node list\nros2 param get /mission_manager stop_distance",
        "capture": "ros2 launch agv_bringup agv_sim.launch.py --show-args; printf '\\n--- config files ---\\n'; find src -path '*/config/*.yaml' -type f | sort",
        "visual": "blocks/F_integration/captures/01_gazebo_final_actual.png", "visual_caption": "실제 Gazebo Sim: launch로 World와 AGV를 함께 실행한 화면",
        "errors": [("일부 노드 시작 실패", "launch 순서·spawn 준비·package path를 확인한다."), ("YAML 값이 적용 안 됨", "node name/namespace와 parameter key를 확인한다.")],
        "mini": "mission.yaml stop_distance를 0.50→0.70으로 바꾼 뒤 재실행하고 ros2 param get으로 적용값을 확인한다.",
        "extra": [("rosbag2 재생 절차", ["기록 종료 뒤 `ros2 bag info ~/agv_bag_01`로 topic과 duration을 먼저 확인한다.", "Gazebo가 없는 terminal에서 processing node를 source·실행한다.", "`ros2 bag play ~/agv_bag_01 --clock`으로 sensor 입력을 다시 공급한다.", "use_sim_time=true 노드는 /clock을 받아 같은 시간축을 사용한다."]), ("증상별 디버깅 도구", ["node 미실행: ros2 node list", "topic/타입/QoS: ros2 topic info --verbose", "parameter: ros2 param get/list", "frame: tf2_tools view_frames와 RViz2", "시뮬레이터: Gazebo GUI와 bridge log"] )],
    },
    {
        "id": "M22", "block": "F", "folder": "blocks/F_integration/M22_final_project",
        "filename": "M22_Final_Project_가상_AGV_자율_미션.pptx", "title": "Final Project — Gazebo 기반 가상 AGV 자율 미션",
        "goal": "모델·주행·bridge·센서·인지·안전·미션·시각화를 하나의 AGV 프로젝트로 연결한다.",
        "completion": "한 launch 실행, RViz/Gazebo 확인, mission log, rosbag 기록을 제출 체크리스트로 남긴다.",
        "start": "M21의 agv_sim.launch.py와 각 기능별 checkpoint가 준비된 상태.",
        "previous": "M21 one-command bringup", "next": "핵심 과정 완료 → SLAM/Nav2 확장",
        "flow": ["START", "SEARCH", "APPROACH", "AVOID / GOAL", "STOP"],
        "files": ["agv_bringup/launch/agv_sim.launch.py", "agv_bringup/config/mission.yaml", "agv_gazebo/worlds/warehouse.sdf", "agv_mission/agv_mission/mission_manager.py", "agv_control/agv_control/safety_controller.py"],
        "actions": [("환경 기준선을 확인한다", "같은 ROS/Gazebo 조합과 workspace overlay가 결과 재현의 시작점이다.", "source /opt/ros/jazzy/setup.bash\ncd ~/ros2_curri/agv_ws\ncolcon build --symlink-install\nsource install/setup.bash"), ("최종 bringup을 실행한다", "최종 launch가 Gazebo·bridge·TF·sensor·mission·safety를 조립한다.", "ros2 launch agv_bringup agv_sim.launch.py"), ("증거 데이터를 남긴다", "화면만이 아니라 topic·log·bag을 함께 남겨야 다른 사람이 검증할 수 있다.", "ros2 node list\nros2 topic list\nros2 bag record -o ~/agv_final_bag /scan /imu/data /odom /camera/image_raw")],
        "run": "cd ~/ros2_curri/agv_ws && source install/setup.bash\nros2 launch agv_bringup agv_sim.launch.py",
        "validate": "ros2 launch agv_bringup agv_sim.launch.py --show-args\nros2 node list\nros2 topic list | rg '/scan|/imu/data|/odom|/camera/image_raw|/cmd_vel'",
        "capture": "colcon list --names-only; ros2 launch agv_bringup agv_sim.launch.py --show-args; rosdep check --from-paths src --ignore-src --rosdistro jazzy --skip-keys ament_python",
        "visual": "blocks/F_integration/captures/02_rviz_final_actual.png", "visual_caption": "실제 RViz2 통합 display: 최종 미션에서 RobotModel·TF·sensor·Path·Marker를 함께 확인하는 화면",
        "errors": [("일부 노드만 실행", "launch 순서·namespace·config 경로를 확인한다."), ("센서/TF 위치 불일치", "header.frame_id와 TF tree를 확인한다."), ("로봇이 불안정", "physics collision·friction·inertia를 M07 순서로 점검한다.")],
        "mini": "완료 증거를 ①Gazebo ②RViz ③ros2 topic ④mission log ⑤rosbag으로 나누어 실제 제출 파일 목록을 작성한다.",
        "extra": [("최종 미션 시나리오", ["1. START: launch와 sensor topic을 확인한다.", "2. SEARCH: target이 없을 때 회전 탐색 command를 관찰한다.", "3. APPROACH: detection 중심 오차가 angular command로 바뀌는지 확인한다.", "4. AVOID/STOP: LiDAR threshold 안에서는 safety가 우선하는지 확인한다.", "5. GOAL/STOP: 목표 도달 뒤 forward /cmd_vel이 0인지 확인한다."]), ("최종 기능 체크표", ["Model/Physics: body·wheel·collision·inertia", "Drive/TF/Odom: /cmd_vel, /odom, frame tree", "Sensors: /camera/image_raw, /scan, /imu/data", "Autonomy: detection, safety priority, mission state", "Operations/Data: launch, YAML, rosbag2, README/checkpoint"]), ("Milestone 1~8 제출 증거", ["1 모델 TF, 2 World spawn, 3 Differential Drive, 4 bridge/RViz", "5 Camera/LiDAR/IMU, 6 safety/perception, 7 FSM/PID, 8 launch/rosbag", "각 milestone마다 command 결과와 정상 화면을 함께 저장한다."]), ("실패를 재현 가능하게 기록한다", ["topic: 이름·타입·publisher/subscriber 수", "TF: frame_id·parent/child·Fixed Frame", "bridge: 방향·Gazebo/ROS type", "physics: collision·friction·inertia", "launch: package share·YAML key·실행 순서"] )],
    },
]


# PDF example policy: do not show a file or a result without telling a beginner
# (1) what it creates, (2) which part to read first, (3) how it is used, and
# (4) where the result is observed.  Specific guidance overrides the generic
# extension-based explanation below.
FILE_GUIDANCE = {
    "agv_description/curriculum_stages/M05/agv.urdf": (
        "몸체·두 바퀴·caster·camera/lidar/imu frame을 link와 joint로 조립한 기본 AGV",
        "각 <link>의 geometry와 각 <joint>의 parent/child를 차례로 작성한다.",
        "check_urdf로 tree를 확인한 뒤 robot_state_publisher + RViz RobotModel에 넣는다.",
        "terminal의 base_link 자식 6개와 RViz의 바퀴·센서 frame"),
    "agv_description/curriculum_stages/M06/agv.urdf.xacro": (
        "반복되는 좌·우 바퀴를 property와 macro 하나로 생성하는 재사용 가능한 AGV 설명",
        "property → macro 정의 → macro 호출 순서로 읽고, wheel_radius 한 값만 바꾼다.",
        "xacro로 URDF를 만든 뒤 check_urdf를 실행한다.",
        "생성 URDF의 left/right wheel과 같은 반지름"),
    "agv_description/curriculum_stages/M07/model.sdf": (
        "Gazebo가 질량·충돌·마찰을 계산할 수 있는 물리 AGV",
        "visual(보임), collision(부딪힘), inertial(무게), friction(접지)을 분리해 작성한다.",
        "gz sdf -k로 문법을 검사하고 값 하나만 바꿔 다시 실행한다.",
        "terminal의 Valid와 물리 태그, Gazebo에서 바닥 위 안정 상태"),
    "agv_description/curriculum_stages/M08/model.config": (
        "Gazebo가 model://agv 이름으로 SDF 모델을 찾게 하는 모델 등록 정보",
        "name과 sdf 파일 이름이 실제 파일과 같은지 먼저 맞춘다.",
        "models/ 아래에 둘 때 Gazebo가 URI를 해석한다.",
        "World의 <uri>model://agv</uri>"),
    "agv_description/curriculum_stages/M08/model.sdf": (
        "World에 배치하기 직전의 AGV 본체·바퀴·joint 모델",
        "M07 물리 요소를 유지한 채 sensor/drive plugin이 아직 없는지 읽는다.",
        "gz sdf -k 뒤 model.config와 함께 Gazebo models 경로에 둔다.",
        "Valid 출력과 World에 spawn된 AGV"),
    "agv_description/curriculum_stages/M08/warehouse.sdf": (
        "바닥·벽·장애물·목표물과 AGV를 배치하는 실습 공간",
        "world → 환경 model → include 순서로 만들고 include의 pose를 마지막에 조정한다.",
        "gz sim으로 World를 연다.",
        "Gazebo Entity Tree의 ground·wall·target·agv"),
    "agv_gazebo/models/agv/model.sdf": (
        "구동·카메라·LiDAR·IMU를 포함해 Gazebo가 시뮬레이션할 완성 AGV",
        "link/joint 다음에 sensor와 DiffDrive plugin을 찾고, topic 이름을 bridge와 비교한다.",
        "Gazebo launch 뒤 각 ROS topic을 확인한다.",
        "Gazebo의 AGV, /camera/image_raw, /scan, /imu/data"),
    "agv_gazebo/worlds/warehouse.sdf": (
        "센서·물리 system plugin과 warehouse 환경을 함께 여는 Gazebo World",
        "Physics/Sensors/Imu plugin → 환경 → include AGV 순서로 읽는다.",
        "World를 다시 열면 Imu system이 /imu/data의 원천을 생성한다.",
        "Gazebo World와 IMU terminal의 frame_id·publish rate"),
    "agv_gazebo/config/bridge.yaml": (
        "Gazebo Transport 메시지와 ROS 2 topic을 방향별로 연결하는 번역표",
        "한 항목에서 ros_topic_name → type → direction을 세 줄씩 비교한다.",
        "parameter_bridge가 이 YAML을 읽도록 launch한다.",
        "ros2 topic info에서 보이는 type과 publisher 수"),
    "agv_vision/agv_vision/yolo_node.py": (
        "카메라 Image를 받아 검출 결과와 debug Image를 내보내는 인지 노드",
        "subscription → cv_bridge 변환 → 검출 → Detection/debug publish 순서로 읽는다.",
        "enable_yolo=false로 fallback부터 실행하고 모델 준비 뒤 true로 바꾼다.",
        "/vision/debug_image의 노란 box와 target 거리"),
    "agv_sensors/agv_sensors/lidar_processor.py": (
        "LaserScan의 전방 sector에서 유효 거리 최소값을 계산하는 안전 입력 노드",
        "ranges 필터 → 각도 index → minimum → /obstacle_distance publish 순서로 읽는다.",
        "launch 뒤 /scan과 /obstacle_distance를 별 터미널에서 본다.",
        "LiDAR terminal의 frame·range와 RViz LaserScan"),
    "agv_sensors/agv_sensors/imu_monitor.py": (
        "IMU frame·각속도·가속도를 단위와 함께 출력하는 진단 노드",
        "message header → angular_velocity → linear_acceleration 순서로 읽는다.",
        "제자리 회전 명령과 함께 실행해 z축 부호를 비교한다.",
        "IMU terminal의 frame_id와 publish rate"),
    "agv_bringup/config/sensors.yaml": (
        "센서 처리 노드가 공통으로 사용할 use_sim_time 등 실행 파라미터",
        "노드 이름 아래 parameter와 값의 단위를 확인한다.",
        "launch가 YAML을 넘긴 뒤 ros2 param get으로 확인한다.",
        "/clock과 use_sim_time=true"),
}


MODULE_FOCUS = {
    "M05": ("URDF는 로봇의 부품표가 아니라, ‘어떤 부품을 어느 frame에 고정할지’ 적는 조립도입니다.",
            "<link name=\"base_link\">, <joint name=\"...\">, <parent>, <child>",
            "wheel link의 cylinder radius를 바꾸면 보이는 바퀴 크기가 바뀌고, joint origin은 장착 위치를 바꿉니다.",
            "check_urdf의 root Link와 child 목록을 먼저 읽고 RViz RobotModel로 형태를 봅니다."),
    "M06": ("Xacro는 같은 바퀴 코드를 복사하지 않고 ‘값 한 번, 생성 여러 번’으로 바꾸는 도구입니다.",
            "<xacro:property>, <xacro:macro>, <xacro:wheel ...>",
            "wheel_radius 값을 하나 바꾸면 좌우 macro 호출에 같은 값이 전달됩니다.",
            "xacro → /tmp/agv.urdf → check_urdf 순서가 변환과 검증을 분리합니다."),
    "M07": ("SDF 물리 모델은 화면에 보이는 모양과 실제 접촉 계산을 서로 다른 태그로 관리합니다.",
            "<visual>, <collision>, <inertial>, <friction>",
            "visual만 바꾸면 모양만 변하고, mass/inertia/friction은 주행·접지 반응을 바꿉니다.",
            "Valid 뒤 Gazebo에서 떠오름·떨림·미끄럼이 없는지 관찰합니다."),
    "M08": ("World는 무대이고 model://agv include는 그 무대 위에 로봇을 놓는 한 줄입니다.",
            "<world>, <include>, <uri>model://agv</uri>, <pose>",
            "include pose x/y/z를 바꾸면 spawn 위치만 바뀌며 AGV 모델 파일은 바뀌지 않습니다.",
            "Gazebo Entity Tree에서 환경 model과 agv가 따로 보이는지 확인합니다."),
    "M12": ("Camera는 SDF에서 영상을 만들고 bridge가 ROS Image로 옮긴 뒤 cv_bridge 코드가 픽셀로 바꿉니다.",
            "<horizontal_fov>, <width>, <height>, <update_rate>, /camera/image_raw", 
            "FOV는 보이는 폭, width/height는 프레임 크기, update_rate는 새 영상의 주기를 바꿉니다.",
            "Camera raw 화면 → ros2 topic info/hz → debug Image 순서로 정상 여부를 확인합니다."),
    "M13": ("LiDAR의 ranges는 거리 배열이고, 각 원소의 방향은 angle_min + index × angle_increment입니다.",
            "<samples>, <min_angle>, <max_angle>, range_min/range_max, qos_profile_sensor_data", 
            "front_half_angle_deg를 바꾸면 최소거리 판단에 포함되는 전방 sector가 바뀝니다.",
            "terminal에서 frame/range/rate를 보고 RViz LaserScan에서 공간 방향을 확인합니다."),
    "M14": ("IMU는 로봇 base에 붙은 좌표계에서 회전과 가속을 계속 발행하는 센서입니다.",
            "<sensor name=\"imu\">, update_rate, Imu system plugin, angular_velocity", 
            "World의 gz-sim-imu-system이 없으면 sensor 태그가 있어도 실제 메시지가 나오지 않습니다.",
            "terminal의 agv/base_link/imu frame과 publish rate를 보고 제자리 회전 시 z축을 비교합니다."),
    "M15": ("센서를 동시에 운영할 때는 ‘topic이 있다’보다 시간·QoS·frame·rate가 서로 맞는지가 중요합니다.",
            "/clock, use_sim_time, qos_profile_sensor_data, ros2 topic hz", 
            "SDF update_rate와 실제 hz가 크게 다르면 시뮬레이터 부하·QoS·bridge를 먼저 점검합니다.",
            "camera/LiDAR/IMU terminal과 RViz를 함께 열어 이름·rate·Fixed Frame을 대조합니다."),
}


# Each path was captured from the live Ubuntu/ROS 2 environment.  The third
# value is the exact observation a learner should make before moving on.
VISUAL_SEQUENCES = {
    "M05": [("blocks/B_robot_build/captures/02_m05_urdf_terminal_actual.png", "실제 terminal: check_urdf가 base_link와 6개 자식 link를 확인한 결과", "root Link=base_link와 wheel·camera·lidar·imu child 이름을 찾습니다.")],
    "M06": [("blocks/B_robot_build/captures/03_m06_xacro_terminal_actual.png", "실제 terminal: Xacro 변환 뒤 생성 URDF를 check_urdf로 검사한 결과", "xacro 명령 다음의 Successfully Parsed XML과 left/right wheel child를 찾습니다.")],
    "M07": [("blocks/B_robot_build/captures/04_m07_sdf_terminal_actual.png", "실제 terminal: SDF Valid 검사와 visual·collision·mass·friction 태그", "Valid 다음에 visual/collision과 mass/friction이 서로 다른 책임임을 읽습니다.")],
    "M08": [("blocks/B_robot_build/captures/05_m08_world_terminal_actual.png", "실제 terminal: World 이름과 model://agv include·spawn pose", "warehouse_m08, model://agv, pose 세 항목을 순서대로 찾습니다."),
            ("blocks/B_robot_build/captures/01_gazebo_world_spawn_actual.png", "실제 Gazebo Sim: warehouse World와 Entity Tree에 보이는 AGV", "오른쪽 Entity Tree에서 ground·wall·target·agv가 따로 생성됐는지 확인합니다.")],
    "M12": [("blocks/D_sensors/captures/01_camera_image_raw_actual.png", "실제 /camera/image_raw: AGV 전방 camera가 target·장애물 방향을 보는 원본 프레임", "빨간 target과 노란 장애물의 화면 위치로 camera pose·FOV를 판단합니다."),
            ("blocks/D_sensors/captures/02_camera_terminal_actual.png", "실제 GNOME Terminal: camera Image type·publisher/subscriber·publish rate", "Image 타입, publisher 1개 이상, average rate를 함께 찾습니다."),
            ("blocks/D_sensors/captures/04_vision_debug_actual.png", "실제 /vision/debug_image: fallback 검출 결과가 target box·거리로 표시된 화면", "노란 box와 ‘target … m’ 텍스트가 raw frame 처리 결과인지 확인합니다.")],
    "M13": [("blocks/D_sensors/captures/05_lidar_terminal_actual.png", "실제 GNOME Terminal: LaserScan frame·angle·range와 publish rate", "frame_id, ±pi angle, range_min/max, average rate를 순서대로 읽습니다."),
            ("blocks/D_sensors/captures/03_rviz_sensors_actual.png", "실제 RViz2: LaserScan을 RobotModel·TF와 같은 Fixed Frame에서 연 화면", "Displays에서 LaserScan이 켜져 있고 scan frame을 TF가 변환하는지 확인합니다.")],
    "M14": [("blocks/D_sensors/captures/06_imu_terminal_actual.png", "실제 GNOME Terminal: /imu/data의 frame·필드와 publish rate", "agv/base_link/imu frame과 angular_velocity·linear_acceleration 필드, average rate를 찾습니다.")],
    "M15": [("blocks/D_sensors/captures/02_camera_terminal_actual.png", "실제 terminal: Camera topic의 type·연결 수·실측 rate", "camera가 message를 내보내고 있는지 type·publisher·rate로 확인합니다."),
            ("blocks/D_sensors/captures/05_lidar_terminal_actual.png", "실제 terminal: LiDAR의 frame·range·실측 rate", "SDF 설정값과 실제 range/angle/rate가 논리적으로 맞는지 대조합니다."),
            ("blocks/D_sensors/captures/06_imu_terminal_actual.png", "실제 terminal: IMU frame·필드·실측 rate", "IMU도 같은 simulation 시간에 갱신되는지 average rate를 확인합니다."),
            ("blocks/D_sensors/captures/03_rviz_sensors_actual.png", "실제 RViz2: sensor display를 같은 Fixed Frame에서 대조", "camera·LiDAR·IMU 각각을 화면·terminal·frame 기준으로 교차 확인합니다.")],
}


def canonical_path(relative: str) -> str:
    return "~/ros2_curri/agv_ws/src/" + relative


def fold_python_line(line: str, width: int = 88) -> list[str]:
    """Fold only after commas inside brackets, preserving valid Python syntax."""
    if len(line) <= width:
        return [line]
    indent = line[:len(line) - len(line.lstrip())]
    continuation = indent + "    "
    result: list[str] = []
    current = indent
    depth = 0
    quote: str | None = None
    escaped = False
    for char in line[len(indent):]:
        if quote:
            current += char
            if escaped:
                escaped = False
            elif char == "\\\\":
                escaped = True
            elif char == quote:
                quote = None
            continue
        if char in {"'", '\"'}:
            quote = char
        elif char in "([{":
            depth += 1
        elif char in ")]}":
            depth = max(0, depth - 1)
        current += char
        if len(current) > width and depth > 0 and "," in current[len(indent):]:
            split_at = current.rfind(",") + 1
            result.append(current[:split_at].rstrip())
            current = continuation + current[split_at:].lstrip()
    result.append(current.rstrip())
    return result


def source_text(relative: str) -> str:
    path = SRC / relative
    text = path.read_text(encoding="utf-8")
    if path.suffix == ".py":
        try:
            # Existing examples keep some packaging files on one long line.
            # ast.unparse preserves the complete program while making the code
            # readable at the guide's required 14 pt code size.
            formatted = ast.unparse(ast.parse(text)).strip()
            return "\n".join(part for line in formatted.splitlines() for part in fold_python_line(line))
        except SyntaxError:
            return text.rstrip()
    if path.suffix in {".sdf", ".xacro", ".xml", ".config"}:
        try:
            formatted = minidom.parseString(text.encode("utf-8")).toprettyxml(indent="  ").strip()
            # XML attributes may legally continue on the next line. Fold only
            # long tag lines so the displayed code remains valid XML/Xacro.
            folded = []
            for line in formatted.splitlines():
                if len(line) <= 88 or not line.lstrip().startswith("<"):
                    folded.append(line)
                    continue
                indent = line[:len(line) - len(line.lstrip())]
                folded.extend(textwrap.wrap(line.strip(), width=88, initial_indent=indent,
                                            subsequent_indent=indent + "  ",
                                            break_long_words=False, break_on_hyphens=False))
            return "\n".join(folded)
        except Exception:
            return text.rstrip()
    return text.rstrip()


def run_actual(command: str) -> str:
    prefix = "export PATH=/usr/bin:/bin:$PATH; source /opt/ros/jazzy/setup.bash; source install/setup.bash; "
    completed = subprocess.run(["bash", "-c", prefix + command], cwd=WORKSPACE, text=True,
                               stdout=subprocess.PIPE, stderr=subprocess.STDOUT, timeout=45, check=False)
    output = completed.stdout.strip()
    return output if output else "(명령이 오류 출력 없이 완료되었습니다.)"


def make_terminal_capture(module: dict, text: str, target: Path) -> None:
    command = module["capture"].replace(str(ROOT), "~/ros2_curri/agv_ws")
    lines = [f"$ cd ~/ros2_curri/agv_ws", f"$ {command}", ""] + text.splitlines()
    wrapped = []
    for line in lines:
        wrapped.extend(textwrap.wrap(line, width=98, replace_whitespace=False, drop_whitespace=False) or [""])
    height = max(430, 54 + min(len(wrapped), 34) * 25)
    image = Image.new("RGB", (1220, height), (25, 31, 40))
    draw = ImageDraw.Draw(image)
    regular = ImageFont.truetype(FONT_FILE, 19)
    draw.rounded_rectangle((0, 0, 1220, 38), radius=0, fill=(47, 56, 70))
    draw.ellipse((16, 12, 28, 24), fill=(237, 86, 85)); draw.ellipse((37, 12, 49, 24), fill=(240, 187, 80)); draw.ellipse((58, 12, 70, 24), fill=(89, 195, 111))
    y = 52
    for line in wrapped[:34]:
        color = (153, 216, 163) if line.startswith("$") else (233, 237, 242)
        draw.text((22, y), line, font=regular, fill=color)
        y += 25
    image.save(target)


def set_text(shape, text: str, size: int = 20, color: RGBColor = NAVY, bold: bool = False,
             font: str = FONT, align=PP_ALIGN.LEFT) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.08)
    frame.margin_top = frame.margin_bottom = Inches(0.04)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.space_after = Pt(0)


def text_box(slide, x, y, w, h, text: str, size: int = 20, color: RGBColor = NAVY,
             bold: bool = False, font: str = FONT, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shape, text, size, color, bold, font, align)
    return shape


def box(slide, x, y, w, h, fill: RGBColor = WHITE, line: RGBColor = RGBColor(212, 220, 230),
        radius=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    return shape


def add_footer(slide, module: dict) -> None:
    text_box(slide, 0.58, 7.12, 8.5, 0.22, f"ROS 2 Jazzy · Gazebo Harmonic · {module['id']} · 따라 하기형 AGV 실습", 10, GREY)
    text_box(slide, 10.0, 7.12, 2.7, 0.22, "명령·파일 경로는 슬라이드에서 복사 가능", 10, GREY, align=PP_ALIGN.RIGHT)


def add_notes(slide, explanation: str, command: str = "", check: str = "", error: str = "") -> None:
    note = (f"[설명] {explanation}\n"
            f"[실행] {command or '슬라이드의 명령 또는 파일 작성 단계를 진행한다.'}\n"
            f"[확인] {check or '슬라이드의 완료 조건을 교육생과 함께 확인한다.'}\n"
            f"[멈춤] 교육생이 현재 체크포인트를 통과할 때까지 기다린다.\n"
            f"[오류] {error or '오류가 나면 현재 폴더, source, 파일 경로를 먼저 확인한다.'}")
    slide.notes_slide.notes_text_frame.text = note


def title(slide, module: dict, text: str, subtitle: str = "") -> None:
    text_box(slide, 0.62, 0.38, 11.9, 0.48, text, 32, NAVY, True)
    if subtitle:
        text_box(slide, 0.66, 0.94, 11.7, 0.32, subtitle, 17, GREY)
    add_footer(slide, module)


def status(slide, module: dict, step: str, folder: str = "~/ros2_curri/agv_ws", new_terminal: str = "현재 터미널") -> None:
    bar = box(slide, 0.62, 1.30, 12.1, 0.38, fill=LIGHT, line=LIGHT, radius=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    set_text(bar, f"현재 단계 {step}    |    현재 폴더 {folder}    |    {new_terminal}", 13, NAVY, True)


def add_bullets(slide, x, y, w, h, bullets: list[str], size: int = 20, color: RGBColor = NAVY) -> None:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame; frame.clear(); frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.12)
    for i, item in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = item; p.font.name = FONT; p.font.size = Pt(size); p.font.color.rgb = color
        p.level = 0; p.text = "• " + p.text; p.space_after = Pt(8)


def add_flow(slide, module: dict) -> None:
    title(slide, module, "핵심 흐름을 한 장으로 이해하기", "상자와 화살표의 방향을 따라 데이터 또는 제어의 흐름을 읽습니다.")
    flow = module["flow"]
    count = len(flow)
    width = min(2.22, 10.8 / count)
    start_x = (13.33 - (count * width + (count - 1) * 0.34)) / 2
    y = 3.0
    for index, label in enumerate(flow):
        x = start_x + index * (width + 0.34)
        shape = box(slide, x, y, width, 1.15, fill=WHITE, line=BLUE)
        set_text(shape, label, 19, NAVY, True, align=PP_ALIGN.CENTER)
        shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        if index < count - 1:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + width), Inches(y + 0.57), Inches(x + width + 0.31), Inches(y + 0.57))
            line.line.color.rgb = BLUE; line.line.width = Pt(2.5); line.line.end_arrowhead = True
    text_box(slide, 1.0, 5.25, 11.3, 0.55, "기억할 한 문장: 앞 단계의 출력이 다음 단계의 입력이 되도록 topic·frame·파일 이름을 끝까지 유지합니다.", 20, BLUE, True, align=PP_ALIGN.CENTER)
    add_notes(slide, "도식의 각 상자가 무엇을 만들고 화살표가 어떤 방향으로 데이터를 보내는지 설명한다.", check="교육생이 왼쪽에서 오른쪽(또는 제어 방향)으로 흐름을 말할 수 있다.")


def file_guidance(relative: str) -> tuple[str, str, str, str]:
    if relative in FILE_GUIDANCE:
        return FILE_GUIDANCE[relative]
    suffix = Path(relative).suffix
    if "/launch/" in relative or relative.endswith("launch.py"):
        return ("여러 ROS 2 node를 시작 순서대로 조립하는 launch 파일", "launch argument → Node → parameters/config 순서로 읽는다.",
                "ros2 launch로 실행한다.", "terminal의 실행 node와 topic")
    if suffix == ".py":
        return ("ROS 2 node 또는 실행 보조 로직", "import → Node 생성 → publisher/subscription → callback/main 순서로 읽는다.",
                "package를 build·source한 뒤 ros2 run 또는 launch에서 실행한다.", "terminal의 node/topic/log 출력")
    if suffix in {".sdf", ".urdf", ".xacro"}:
        return ("로봇 또는 Gazebo 환경의 구조·물리·sensor 설정", "link/model → geometry → joint/sensor/plugin 순서로 읽는다.",
                "문법 검사 뒤 Gazebo/RViz에서 열어 본다.", "terminal 검사 결과와 3D 화면")
    if suffix in {".yaml", ".yml"}:
        return ("실행 시 읽히는 ROS 2 설정값", "이름 → type/value → launch에서 전달되는 위치를 비교한다.",
                "launch 실행 뒤 ros2 param/topic 명령으로 확인한다.", "terminal의 parameter/topic 출력")
    return ("이번 모듈의 재현 가능한 파일", "파일명·경로·핵심 이름을 슬라이드와 동일하게 작성한다.",
            "build/source 뒤 실행 단계에서 사용한다.", "검증 명령의 파일·topic 이름")


def add_file_role_slide(slide, module: dict, group: list[str], page_no: int, page_count: int) -> None:
    title(slide, module, f"코드/설정이 실제로 만드는 것 ({page_no}/{page_count})", "파일을 입력하기 전에 ‘무엇을 만들고 어디서 볼지’를 먼저 연결합니다.")
    status(slide, module, "파일 역할", "~/ros2_curri/agv_ws")
    for index, relative in enumerate(group):
        makes, read_order, use, observe = file_guidance(relative)
        y = 1.78 + index * 2.53
        panel = box(slide, 0.70, y, 12.0, 2.24, fill=LIGHT, line=RGBColor(204, 219, 235))
        text_box(slide, 0.93, y + 0.15, 11.4, 0.28, canonical_path(relative), 13, BLUE, True, MONO)
        text_box(slide, 0.95, y + 0.55, 1.35, 0.25, "만드는 것", 15, GREEN, True)
        text_box(slide, 2.16, y + 0.53, 10.05, 0.30, makes, 15, NAVY)
        text_box(slide, 0.95, y + 0.96, 1.35, 0.25, "코드 읽기", 15, ORANGE, True)
        text_box(slide, 2.16, y + 0.94, 10.05, 0.33, read_order, 14, NAVY)
        text_box(slide, 0.95, y + 1.38, 1.35, 0.25, "사용/확인", 15, BLUE, True)
        text_box(slide, 2.16, y + 1.36, 10.05, 0.56, f"사용: {use}\n화면/출력: {observe}", 13, NAVY)
    add_notes(slide, "파일을 코드 문법보다 먼저 역할·사용법·확인 화면으로 설명한다.", check="교육생이 각 파일의 산출물과 확인 위치를 말할 수 있다.")


def add_code_reading_slide(slide, module: dict, focus: tuple[str, str, str, str]) -> None:
    purpose, markers, change, observe = focus
    title(slide, module, "코드를 ‘만드는 것 → 표시할 줄 → 바꾼 결과’로 읽는다", "PDF 예시처럼 화면을 보기 전에 코드의 역할과 관찰 지점을 한 번 짚고 갑니다.")
    status(slide, module, "코드 해설", "~/ros2_curri/agv_ws")
    panels = [
        ("1. 이 파일이 만드는 것", purpose, GREEN),
        ("2. 먼저 찾을 코드/태그", markers, BLUE),
        ("3. 값을 바꾸면", change, ORANGE),
        ("4. 실제로 보는 곳", observe, RED),
    ]
    for index, (heading, body, color) in enumerate(panels):
        x = 0.72 + (index % 2) * 6.08
        y = 1.82 + (index // 2) * 2.35
        panel = box(slide, x, y, 5.82, 1.93, fill=LIGHT, line=color)
        text_box(slide, x + 0.22, y + 0.18, 5.35, 0.29, heading, 17, color, True)
        text_box(slide, x + 0.22, y + 0.63, 5.30, 1.03, body, 16, NAVY)
    text_box(slide, 0.85, 6.75, 11.6, 0.22, "다음 코드 슬라이드에서는 위의 태그/함수를 찾아 색으로 표시하며 한 줄씩 설명합니다.", 12, GREEN, True, align=PP_ALIGN.CENTER)
    add_notes(slide, purpose, check=observe)


def add_code_slide(slide, module: dict, relative: str, lines: list[str], chunk_no: int, chunk_count: int, step_no: int) -> None:
    display_path = canonical_path(relative)
    title(slide, module, f"파일을 {'새로 만들고' if chunk_no == 1 else '계속 작성하고'} 저장한다 ({chunk_no}/{chunk_count})", display_path)
    status(slide, module, f"{step_no}/{step_no + chunk_count - chunk_no + 1}", "~/ros2_curri/agv_ws")
    text_box(slide, 0.72, 1.78, 12.0, 0.28, "이번 단계: 파일의 전체 내용을 중간 생략 없이 작성합니다. 아래 줄 범위를 끝까지 입력하세요.", 15, NAVY, True)
    command_box = box(slide, 0.72, 2.14, 12.0, 0.66, fill=LIGHT, line=LIGHT, radius=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    directory = str(Path(display_path).parent)
    set_text(command_box, f"파일 생성: mkdir -p {directory}\n파일 열기: nano {display_path}", 14, NAVY, False, MONO)
    code_box = box(slide, 0.72, 2.96, 12.0, 3.68, fill=DARK, line=DARK, radius=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    code_shape = code_box
    code_text = "\n".join(lines)
    set_text(code_shape, code_text, 14, WHITE, False, MONO)
    code_shape.text_frame.margin_left = code_shape.text_frame.margin_right = Inches(0.18)
    code_shape.text_frame.margin_top = Inches(0.12)
    text_box(slide, 0.82, 6.76, 11.8, 0.22, f"파일 전체 {chunk_no}/{chunk_count} · 이 슬라이드 다음에 같은 파일의 나머지 줄이 이어집니다.", 10, GREY)
    add_notes(slide, f"{display_path}의 전체 파일을 {chunk_no}/{chunk_count} 구간으로 작성한다.", command=f"nano {display_path}", check="코드를 저장하고 다음 코드 슬라이드 또는 build 단계로 이동한다.", error="들여쓰기·XML 태그·따옴표가 슬라이드와 정확히 같은지 확인한다.")


def add_action_slide(slide, module: dict, step_number: int, total_steps: int, action: tuple[str, str, str]) -> None:
    heading, why, command = action
    title(slide, module, heading, "한 장에는 한 행동만 수행합니다.")
    status(slide, module, f"{step_number}/{total_steps}", "~/ros2_curri/agv_ws", "새 터미널" if "새 터미널" in command else "현재 터미널")
    text_box(slide, 0.74, 1.93, 2.0, 0.38, "왜 하는가", 18, BLUE, True)
    text_box(slide, 0.74, 2.30, 11.8, 0.56, why, 20, NAVY)
    command_box = box(slide, 0.74, 3.14, 11.85, 2.55, fill=DARK, line=DARK, radius=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE)
    set_text(command_box, command, 16, WHITE, False, MONO)
    text_box(slide, 0.84, 5.98, 1.25, 0.28, "확인", 18, GREEN, True)
    text_box(slide, 0.84, 6.28, 11.55, 0.38, "명령의 출력에서 topic·frame·파일 이름 또는 성공 메시지를 찾으면 이 단계가 완료됩니다.", 18, NAVY)
    add_notes(slide, why, command=command, check="명령 출력에서 슬라이드의 확인 항목을 찾는다.")


def add_capture_slide(slide, module: dict, image_path: Path, caption: str, observation: str = "") -> None:
    title(slide, module, "정상 결과: 실제 환경에서 확인한 출력", caption)
    with Image.open(image_path) as image:
        ratio = image.width / image.height
    max_w, max_h = 11.6, 5.25
    width = min(max_w, max_h * ratio); height = width / ratio
    slide.shapes.add_picture(str(image_path), Inches((13.33 - width) / 2), Inches(1.48 + (max_h - height) / 2), width=Inches(width), height=Inches(height))
    read_line = observation or "① 명령/설정  ② 정상 출력 또는 화면  ③ 다음 검증 명령으로 재확인"
    text_box(slide, 0.85, 6.72, 11.6, 0.34, f"화면에서 확인: {read_line}", 11, GREEN, True, align=PP_ALIGN.CENTER)
    add_notes(slide, "실제 Ubuntu/ROS 2 환경에서 생성한 검증 화면이다.", check=read_line)


def add_validation_slide(slide, module: dict) -> None:
    title(slide, module, "객관적 검증 명령으로 완료를 확인한다", "화면만 보지 말고 topic·frame·파일·parameter를 다시 확인합니다.")
    status(slide, module, "검증", "~/ros2_curri/agv_ws")
    command_box = box(slide, 0.72, 1.92, 12.0, 2.35, fill=DARK, line=DARK)
    set_text(command_box, module["validate"], 16, WHITE, False, MONO)
    check_box = box(slide, 0.72, 4.60, 12.0, 1.70, fill=LIGHT, line=RGBColor(209, 223, 213))
    set_text(check_box, "완료 체크\n□ 명령이 오류 없이 실행된다.\n□ 기대한 파일·topic·frame·parameter 이름이 보인다.\n□ 값의 단위와 방향을 한 문장으로 설명할 수 있다.", 18, GREEN, True)
    text_box(slide, 0.80, 6.55, 11.8, 0.26, f"완료 조건: {module['completion']}", 15, NAVY, True)
    add_notes(slide, "검증 명령은 GUI 결과와 별도로 실행해 재현 가능한 완료 기준을 만든다.", command=module["validate"], check=module["completion"])


def add_errors_slide(slide, module: dict) -> None:
    title(slide, module, "막히면: 증상 → 우선 확인 → 재검증", "오류 메시지를 숨기지 않고 같은 순서로 점검합니다.")
    status(slide, module, "오류 대응", "~/ros2_curri/agv_ws")
    errors = module["errors"]
    for index, (symptom, check) in enumerate(errors):
        y = 1.92 + index * 1.55
        panel = box(slide, 0.76, y, 12.0, 1.22, fill=RGBColor(255, 247, 244), line=RGBColor(237, 190, 180))
        text_box(slide, 0.98, y + 0.16, 3.2, 0.30, f"증상 {index + 1}: {symptom}", 18, RED, True)
        text_box(slide, 4.15, y + 0.14, 8.25, 0.68, f"우선 확인: {check}\n수정 후: {module['validate'].splitlines()[0]}", 16, NAVY)
    if len(errors) <= 2:
        text_box(slide, 0.9, 5.58, 11.6, 0.44, "공통 순서: ① 현재 폴더 ② source ③ 파일 경로 ④ 이름/타입/단위 ⑤ build·source 후 재실행", 18, ORANGE, True, align=PP_ALIGN.CENTER)
    else:
        text_box(slide, 0.9, 6.32, 11.6, 0.54, "공통 순서: ① 현재 폴더 ② source ③ 파일 경로\n④ 이름/타입/단위 ⑤ build·source 후 재실행", 15, ORANGE, True, align=PP_ALIGN.CENTER)
    add_notes(slide, "오류를 대신 해결하지 말고 증상부터 우선 확인 순서대로 교육생이 실행하게 한다.", check="수정 뒤 첫 검증 명령을 다시 실행한다.", error="슬라이드의 두 대표 오류를 먼저 사용한다.")


def add_checkpoint_slide(slide, module: dict, artifact_dir: Path) -> None:
    title(slide, module, "체크포인트를 저장하고 다음 모듈로 넘긴다", "현재 Complete는 다음 Starter의 기준선입니다.")
    status(slide, module, "체크포인트", "~/ros2_curri/agv_ws")
    panel = box(slide, 0.75, 1.90, 12.0, 3.95, fill=LIGHT, line=RGBColor(210, 220, 232))
    items = [
        f"□ 완료 화면: {module['completion']}",
        f"□ Complete 파일: {artifact_dir.relative_to(ROOT)}/complete/",
        f"□ 정상 로그·캡처: logs/validation.log, screenshots/validation_terminal.png",
        "□ 버전 식별: CHECKSUM_or_TAG.txt의 SHA-256 manifest",
        f"□ 다음 시작 조건: {module['next']}",
    ]
    add_bullets(slide, 0.98, 2.18, 11.3, 3.25, items, 18)
    text_box(slide, 0.95, 6.12, 11.5, 0.37, "다음 모듈을 시작하기 전, 이 슬라이드의 Complete와 다음 모듈의 Starter를 비교합니다.", 17, GREEN, True, align=PP_ALIGN.CENTER)
    add_notes(slide, "교육생이 Complete·캡처·로그를 저장했는지 확인하고 다음 모듈의 시작 상태를 읽는다.", check=f"{module['next']} 시작 조건을 말할 수 있다.")


def add_reference_slide(slide, module: dict) -> None:
    title(slide, module, "공식 참고 자료", "본문 단계는 링크를 열지 않아도 완료할 수 있으며, 링크는 개념을 더 확인할 때만 사용합니다.")
    refs = [
        ("ROS 2 Jazzy — Gazebo Simulation", "https://docs.ros.org/en/jazzy/Tutorials/Advanced/Simulators/Gazebo/Simulation-Gazebo.html"),
        ("ROS 2 Jazzy — URDF / robot_state_publisher", "https://docs.ros.org/en/jazzy/Tutorials/Intermediate/URDF/Using-URDF-with-Robot-State-Publisher-py.html"),
        ("Gazebo — ROS 2 Integration / ros_gz_bridge", "https://gazebosim.org/docs/latest/ros2_integration/"),
        ("ROS 2 Jazzy — rosbag2 Recording and Playback", "https://docs.ros.org/en/jazzy/Tutorials/Beginner-CLI-Tools/Recording-And-Playing-Back-Data/Recording-And-Playing-Back-Data.html"),
    ]
    panel = box(slide, 0.72, 1.70, 12.0, 4.95, fill=LIGHT, line=RGBColor(215, 223, 233))
    for index, (name, url) in enumerate(refs):
        y = 1.94 + index * 1.18
        text_box(slide, 1.0, y, 4.2, 0.28, name, 17, NAVY, True)
        text_box(slide, 1.0, y + 0.30, 10.9, 0.52, url, 10, BLUE, False, MONO)
    add_notes(slide, "공식 문서는 버전 차이와 확장 학습을 확인할 때 사용한다.", check="학습자는 본문만으로 현재 모듈을 완료한 상태여야 한다.")


def add_extra_slide(slide, module: dict, heading: str, bullets: list[str]) -> None:
    title(slide, module, heading, "통합 단계에서 놓치기 쉬운 실행·운영 조건을 정리합니다.")
    status(slide, module, "통합 보강", "~/ros2_curri/agv_ws")
    panel = box(slide, 0.74, 1.85, 12.0, 4.85, fill=LIGHT, line=RGBColor(210, 222, 236))
    add_bullets(slide, 1.0, 2.18, 11.3, 4.10, bullets, 19)
    add_notes(slide, heading, check="각 항목이 현재 시스템의 어느 package/topic/file과 연결되는지 확인한다.")


def build_deck(module: dict) -> tuple[Path, str]:
    folder = ROOT / module["folder"]
    folder.mkdir(parents=True, exist_ok=True)
    capture_dir = folder / "screenshots"; capture_dir.mkdir(exist_ok=True)
    log_dir = folder / "logs"; log_dir.mkdir(exist_ok=True)
    output_log = run_actual(module["capture"])
    (log_dir / "validation.log").write_text(output_log + "\n", encoding="utf-8")
    terminal_image = capture_dir / "validation_terminal.png"
    make_terminal_capture(module, output_log, terminal_image)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    # 1. Cover
    slide = presentation.slides.add_slide(blank)
    banner = box(slide, 0, 0, 13.333, 1.15, fill=NAVY, line=NAVY, radius=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    set_text(banner, f"{module['id']} · Block {module['block']} · ROS 2 + Gazebo Sim AGV", 18, WHITE, True, align=PP_ALIGN.CENTER)
    text_box(slide, 0.78, 1.75, 11.7, 1.1, module["title"], 34, NAVY, True, align=PP_ALIGN.CENTER)
    text_box(slide, 1.1, 3.10, 11.1, 0.55, module["completion"], 22, GREEN, True, align=PP_ALIGN.CENTER)
    panel = box(slide, 1.5, 4.18, 10.3, 1.15, fill=LIGHT, line=LIGHT)
    set_text(panel, "기준 환경\nUbuntu 24.04 LTS · ROS 2 Jazzy · Gazebo Harmonic · RViz2 · Python(rclpy)", 18, NAVY, True, align=PP_ALIGN.CENTER)
    text_box(slide, 0.7, 6.65, 12.0, 0.28, "이 PPT만 보고 시작 상태 → 파일 작성 → build/run → 검증 → 오류 대응 → checkpoint를 순서대로 수행합니다.", 15, GREY, align=PP_ALIGN.CENTER)
    add_footer(slide, module)
    add_notes(slide, module["goal"], check=module["completion"])

    # 2. roadmap
    slide = presentation.slides.add_slide(blank)
    title(slide, module, "전체 과정에서 현재 위치를 확인한다", "앞 모듈의 Complete가 현재 Starter가 되고, 현재 Complete가 다음 모듈로 넘어갑니다.")
    previous, current, next_module = module["previous"], module["id"] + "\n" + module["title"], module["next"]
    labels = [(previous, GREY), (current, BLUE), (next_module, GREEN)]
    for index, (label, color) in enumerate(labels):
        x = 0.78 + index * 4.2
        shape = box(slide, x, 2.60, 3.55, 1.45, fill=WHITE, line=color)
        set_text(shape, label, 20, color, True, align=PP_ALIGN.CENTER); shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        if index < 2:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 3.56), Inches(3.32), Inches(x + 4.06), Inches(3.32))
            line.line.color.rgb = color; line.line.width = Pt(3); line.line.end_arrowhead = True
    text_box(slide, 0.95, 5.05, 11.4, 0.62, f"이번 모듈의 인계: {module['completion']}", 22, NAVY, True, align=PP_ALIGN.CENTER)
    add_notes(slide, "현재 모듈이 앞뒤 과정과 분리된 예제가 아님을 보여 준다.", check=f"이전={module['previous']}, 다음={module['next']}을 말할 수 있다.")

    # 3. start
    slide = presentation.slides.add_slide(blank)
    title(slide, module, "시작 상태를 먼저 맞춘다", "실습 시작 전에 현재 폴더·이전 파일·실행 프로그램을 확인합니다.")
    status(slide, module, "시작", "~/ros2_curri/agv_ws")
    panel = box(slide, 0.78, 1.95, 12.0, 3.75, fill=LIGHT, line=RGBColor(215, 225, 235))
    add_bullets(slide, 1.0, 2.20, 11.4, 3.10, [f"시작 상태: {module['start']}", f"이전 Complete: {module['previous']}", "현재 폴더: ~/ros2_curri/agv_ws", "새 터미널마다: source /opt/ros/jazzy/setup.bash → source install/setup.bash", "파일 경로는 PPT와 Complete 폴더에서 같은 이름으로 유지"], 19)
    text_box(slide, 0.98, 6.05, 11.5, 0.35, "확인: pwd가 ~/ros2_curri/agv_ws이고, 이전 모듈의 완료 조건을 한 문장으로 설명할 수 있으면 시작합니다.", 16, GREEN, True, align=PP_ALIGN.CENTER)
    add_notes(slide, "교육생의 시작 상태를 통일한다.", command="pwd; source /opt/ros/jazzy/setup.bash; source ~/ros2_curri/agv_ws/install/setup.bash", check="현재 폴더와 source 상태를 모두 확인한다.")

    # 4. goal
    slide = presentation.slides.add_slide(blank)
    title(slide, module, "이번 모듈의 눈으로 확인할 완료 조건", "‘이해했다’가 아니라 실제로 보이는 결과를 기준으로 합니다.")
    goal_box = box(slide, 0.85, 1.86, 11.75, 1.45, fill=RGBColor(238, 247, 241), line=RGBColor(168, 211, 181))
    set_text(goal_box, module["goal"], 25, NAVY, True, align=PP_ALIGN.CENTER); goal_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    done_box = box(slide, 0.85, 3.72, 11.75, 1.45, fill=LIGHT, line=RGBColor(214, 224, 235))
    set_text(done_box, "완료 화면/결과\n" + module["completion"], 22, GREEN, True, align=PP_ALIGN.CENTER); done_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_box(slide, 0.95, 5.85, 11.4, 0.45, "종료 전에는 ‘무엇이 보이면 성공인가’를 terminal 또는 Gazebo/RViz 화면에서 반드시 확인합니다.", 18, NAVY, align=PP_ALIGN.CENTER)
    add_notes(slide, module["goal"], check=module["completion"])

    # 5 concept
    slide = presentation.slides.add_slide(blank); add_flow(slide, module)

    # 6 file tree
    slide = presentation.slides.add_slide(blank)
    title(slide, module, "이번 모듈에서 생성·수정하는 파일", "파일명과 package 경로를 끝까지 동일하게 사용합니다.")
    files = module["files"] or ["(M01은 사용자 파일 없이 ROS 2 내장 demo node를 사용합니다.)"]
    tree_lines = ["~/ros2_curri/agv_ws/src/"] + ["└── " + value for value in files]
    panel = box(slide, 0.85, 1.75, 11.8, 4.82, fill=DARK, line=DARK)
    set_text(panel, "\n".join(tree_lines), 16, WHITE, False, MONO)
    text_box(slide, 0.96, 6.72, 11.3, 0.22, "파일 전체 코드는 다음 슬라이드에 이어집니다. 중간 줄을 빼지 않고 끝까지 제공합니다.", 11, GREY, align=PP_ALIGN.CENTER)
    add_notes(slide, "실제 파일 경로를 보여 주고 어떤 package의 책임인지 설명한다.", check="교육생이 파일의 부모 package를 찾을 수 있다.")

    # Before the complete source appears, connect every file to the robot part,
    # running command, and visible result it is responsible for.
    if module["files"]:
        role_groups = [module["files"][index:index + 2] for index in range(0, len(module["files"]), 2)]
        for index, group in enumerate(role_groups, start=1):
            slide = presentation.slides.add_slide(blank)
            add_file_role_slide(slide, module, group, index, len(role_groups))
    if module["id"] in MODULE_FOCUS:
        slide = presentation.slides.add_slide(blank)
        add_code_reading_slide(slide, module, MODULE_FOCUS[module["id"]])

    # action slides before code
    action_count = len(module["actions"])
    for index, action in enumerate(module["actions"], start=1):
        slide = presentation.slides.add_slide(blank)
        add_action_slide(slide, module, index, action_count, action)

    # Complete source code, split in readable 14-line slides.
    code_step = action_count + 1
    for relative in module["files"]:
        lines = source_text(relative).splitlines() or [""]
        # A code line can wrap once on a projector. Keep each code pane below
        # its visual height limit instead of packing in the theoretical 14 rows.
        chunks = [lines[index:index + 12] for index in range(0, len(lines), 12)]
        for number, chunk in enumerate(chunks, start=1):
            slide = presentation.slides.add_slide(blank)
            add_code_slide(slide, module, relative, chunk, number, len(chunks), code_step)
            code_step += 1

    # build/run, real capture, optional GUI capture, validation, errors, mini, extras, checkpoint, references.
    slide = presentation.slides.add_slide(blank)
    title(slide, module, "build → source → run 순서로 실행한다", "새 파일 또는 수정 파일은 build와 source 뒤에만 실행 이름·설정에 반영됩니다.")
    status(slide, module, "실행", "~/ros2_curri/agv_ws")
    panel = box(slide, 0.72, 1.86, 12.0, 3.48, fill=DARK, line=DARK)
    set_text(panel, module["run"], 16, WHITE, False, MONO)
    text_box(slide, 0.86, 5.73, 1.38, 0.26, "확인 기준", 18, GREEN, True)
    text_box(slide, 0.86, 6.05, 11.45, 0.46, module["completion"], 18, NAVY)
    add_notes(slide, "build/source/run을 생략 없이 실행한다.", command=module["run"], check=module["completion"])

    slide = presentation.slides.add_slide(blank); add_capture_slide(slide, module, terminal_image, "실제 ROS 2 환경의 파일·패키지·구성 검증 로그")
    visual_items = VISUAL_SEQUENCES.get(module["id"])
    if visual_items is None and module.get("visual"):
        visual_items = [(module["visual"], module["visual_caption"], "명령·설정과 화면의 topic·frame·결과 이름을 대조합니다.")]
    for visual_relative, visual_caption, observation in visual_items or []:
        visual = ROOT / visual_relative
        if visual.exists():
            shutil.copy2(visual, capture_dir / visual.name)
            slide = presentation.slides.add_slide(blank)
            add_capture_slide(slide, module, visual, visual_caption, observation)
    slide = presentation.slides.add_slide(blank); add_validation_slide(slide, module)
    slide = presentation.slides.add_slide(blank); add_errors_slide(slide, module)
    slide = presentation.slides.add_slide(blank)
    title(slide, module, "미니 실습: 값 하나를 바꾸고 결과를 비교한다", "정상 결과를 만든 뒤에는 한 번에 하나의 값만 바꿉니다.")
    status(slide, module, "미니 실습", "~/ros2_curri/agv_ws")
    panel = box(slide, 0.82, 2.05, 11.65, 2.18, fill=RGBColor(255, 248, 235), line=RGBColor(236, 198, 128))
    set_text(panel, module["mini"], 23, NAVY, True, align=PP_ALIGN.CENTER); panel.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_box(slide, 0.96, 5.02, 11.2, 0.65, "기록할 것: 바꾼 값(단위) · 기대한 변화 · 실제 출력/화면 · 원래 값으로 되돌린 결과", 19, ORANGE, True, align=PP_ALIGN.CENTER)
    add_notes(slide, "교육생이 직접 한 값만 바꿔 원인과 결과를 연결하게 한다.", check="원래 값과 변경 값을 모두 기록한다.")
    for heading, bullets in module.get("extra", []):
        slide = presentation.slides.add_slide(blank); add_extra_slide(slide, module, heading, bullets)
    slide = presentation.slides.add_slide(blank); add_checkpoint_slide(slide, module, folder)
    slide = presentation.slides.add_slide(blank); add_reference_slide(slide, module)

    output = folder / module["filename"]
    presentation.save(output)
    return output, output_log


def make_artifacts(module: dict, validation_log: str) -> None:
    folder = ROOT / module["folder"]
    starter = folder / "starter"; complete = folder / "complete"
    starter.mkdir(exist_ok=True); complete.mkdir(exist_ok=True)
    (starter / "README.md").write_text(
        f"# {module['id']} Starter\n\n"
        f"- 시작 상태: {module['start']}\n"
        f"- 이전 모듈 기준: {module['previous']}\n"
        "- 이 폴더의 PPT를 1번 슬라이드부터 순서대로 수행한다.\n"
        "- 실제 작업 workspace는 `~/ros2_curri/agv_ws`이고, source 원본은 이 저장소의 `agv_ws/src/`이다.\n",
        encoding="utf-8")
    manifest = []
    for relative in module["files"]:
        source = SRC / relative
        destination = complete / "agv_ws" / "src" / relative
        destination.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(source, destination)
        digest = hashlib.sha256(source.read_bytes()).hexdigest()
        manifest.append(f"{digest}  agv_ws/src/{relative}")
    (complete / "README.md").write_text(
        f"# {module['id']} Complete\n\n"
        f"PPT 완료 후의 핵심 파일 snapshot입니다. 완료 조건: {module['completion']}\n\n"
        "파일의 SHA-256 값은 상위 `CHECKSUM_or_TAG.txt`에 기록했습니다.\n",
        encoding="utf-8")
    (folder / "CHECKSUM_or_TAG.txt").write_text(
        f"{module['id']} complete checkpoint\n"
        "SHA-256 manifest (source snapshot at PPT generation time)\n" + ("\n".join(manifest) or "No user source file: ROS 2 built-in demo module.") + "\n",
        encoding="utf-8")
    source_revision = os.environ.get("COURSE_SOURCE_COMMIT", "uncommitted-source")
    (folder / "screenshots" / "SOURCE_COMMIT.txt").write_text(
        f"Source revision: {source_revision}\n"
        "This module's validation terminal and ordered GUI captures were generated against this source revision.\n",
        encoding="utf-8")
    # This log is intentionally copied once more after build_deck so it is present
    # even when a module has no complete source file.
    (folder / "logs" / "validation.log").write_text(validation_log + "\n", encoding="utf-8")


def main() -> None:
    built = []
    for module in MODULES:
        print(f"Creating {module['id']} ...", flush=True)
        output, log = build_deck(module)
        make_artifacts(module, log)
        built.append(output)
    print(f"Created {len(built)} module decks")
    for path in built:
        print(path.relative_to(ROOT))


if __name__ == "__main__":
    main()
