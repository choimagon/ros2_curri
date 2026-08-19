#!/usr/bin/env python3
"""Create Korean Block presentations with real terminal and GUI captures."""
from __future__ import annotations

import os
import subprocess
import textwrap
from datetime import datetime
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
WORKSPACE = ROOT / "agv_ws"
FONT = "Noto Sans CJK KR"
FONT_FILE = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
FONT_BOLD_FILE = "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc"


def shell(command: str) -> str:
    """Run the command on the actual ROS installation and return its terminal text."""
    env_prefix = "export PATH=/usr/bin:/bin:$PATH; source /opt/ros/jazzy/setup.bash; source install/setup.bash; "
    completed = subprocess.run(
        ["bash", "-c", env_prefix + command],
        cwd=WORKSPACE,
        text=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.STDOUT,
        timeout=45,
        check=False,
    )
    text = completed.stdout.strip()
    return text if text else "(명령은 오류 출력 없이 완료되었습니다.)"


BLOCKS = [
    {
        "key": "A",
        "folder": "blocks/A_ros2_basics",
        "title": "Block A — ROS 2 기초",
        "modules": "M01 ROS 2 개념 · M02 Python/C++ Workspace · M03 Pub/Sub · M04 TF2",
        "goal": "ROS 2 노드·토픽 흐름을 확인하고 AGV의 기본 TF tree를 이해한다.",
        "files": [
            "agv_control/agv_control/counter_publisher.py",
            "agv_control/agv_control/counter_monitor.py",
            "agv_cpp_examples/src/status_publisher.cpp",
            "agv_cpp_examples/CMakeLists.txt",
            "agv_description/urdf/agv.urdf.xacro",
            "agv_description/launch/display.launch.py",
        ],
        "commands": """source /opt/ros/jazzy/setup.bash
cd /home/lab4090/ros2_curri/agv_ws
colcon build --symlink-install --packages-select agv_control agv_description agv_cpp_examples
source install/setup.bash
ros2 run agv_control counter_publisher
# 새 터미널: ros2 run agv_control counter_monitor
# C++ 비교: ros2 run agv_cpp_examples status_publisher""",
        "command_notes": [
            "source: 현재 터미널에 ROS 2 Jazzy와 workspace overlay를 등록한다.",
            "ros2 run: setup.py에 등록된 console script를 찾아 독립 노드로 실행한다.",
            "counter_publisher는 1초 timer·queue depth 10으로 /counter를 발행하며, 속도 명령은 -p linear_speed:=값으로 조절한다.",
        ],
        "capture_command": """tmp=$(mktemp -d /tmp/agv_ppt_a.XXXXXX); export ROS_LOG_DIR=$tmp; timeout --signal=INT 4s ros2 run agv_cpp_examples status_publisher --ros-args -p message:='C++ 노드 시작' >$tmp/publisher.log 2>&1 & pub=$!; sleep 1; timeout --signal=INT 2s ros2 topic echo /cpp_status --once 2>&1 || true; wait $pub || true; cat $tmp/publisher.log""",
        "checks": "C++ publisher의 /cpp_status와 Python counter_monitor의 /counter가 각각 수신되면 두 언어의 ROS 2 node 구조가 정상이다. 다음으로 M04에서 RViz Fixed Frame과 TF tree를 확인한다.",
    },
    {
        "key": "B",
        "folder": "blocks/B_robot_build",
        "title": "Block B — AGV 로봇 제작",
        "modules": "M05 URDF · M06 Xacro · M07 물리 · M08 Gazebo World",
        "goal": "body, 바퀴, caster, LiDAR, camera, IMU를 가진 AGV 모델을 URDF/Xacro와 SDF로 만든다.",
        "files": [
            "agv_description/urdf/common.xacro",
            "agv_description/urdf/wheels.xacro",
            "agv_description/urdf/sensors.xacro",
            "agv_gazebo/models/agv/model.sdf",
            "agv_gazebo/worlds/warehouse.sdf",
        ],
        "commands": """cd /home/lab4090/ros2_curri/agv_ws
source /opt/ros/jazzy/setup.bash
source install/setup.bash
xacro src/agv_description/urdf/agv.urdf.xacro > /tmp/agv.urdf
check_urdf /tmp/agv.urdf
ros2 launch agv_gazebo gazebo.launch.py""",
        "command_notes": [
            "xacro: macro/property/include를 순수 URDF로 펼친다. 원본의 body·wheel 치수를 바꾸면 생성 결과가 함께 바뀐다.",
            "check_urdf: root link와 joint 연결을 검사하지만 Gazebo 물리를 실행하지는 않는다.",
            "ros2 launch: model URI 경로·World·bridge를 한 LaunchDescription으로 시작한다.",
        ],
        "visual_capture": "captures/gazebo_agv_actual.png",
        "visual_caption": "실제 Gazebo Sim: warehouse World에 spawn된 agv entity",
        "capture_command": """xacro src/agv_description/urdf/agv.urdf.xacro > /tmp/agv_ppt.urdf && check_urdf /tmp/agv_ppt.urdf && gz sdf -k src/agv_gazebo/models/agv/model.sdf""",
        "checks": "check_urdf 출력에 base_footprint → base_link와 두 바퀴·세 센서 link가 보이면 모델 구조가 정상이다. Gazebo에서 바닥 관통이나 진동이 없는지 이어서 확인한다.",
    },
    {
        "key": "C",
        "folder": "blocks/C_drive_visualization",
        "title": "Block C — 주행과 시각화",
        "modules": "M09 Differential Drive · M10 ros_gz_bridge · M11 RViz2",
        "goal": "Gazebo의 DiffDrive에 /cmd_vel을 전달하고 odom·scan을 ROS/RViz에서 확인한다.",
        "files": [
            "agv_gazebo/models/agv/model.sdf",
            "agv_gazebo/config/bridge.yaml",
            "agv_gazebo/launch/gazebo.launch.py",
            "agv_control/agv_control/cmd_test_node.py",
        ],
        "commands": """ros2 launch agv_gazebo gazebo.launch.py
# 새 터미널
ros2 topic pub --rate 10 /cmd_vel geometry_msgs/msg/Twist \\
  \"{linear: {x: 0.15}, angular: {z: 0.0}}\"
ros2 topic echo /odom --once
ros2 topic echo /scan --once
rviz2""",
        "command_notes": [
            "ros2 topic pub --rate 10: Twist를 초당 10회 보내 DiffDrive가 연속 주행하도록 한다.",
            "linear.x(m/s)와 angular.z(rad/s)를 조절해 직진·회전·곡선 주행을 만든다.",
            "ros_gz_bridge는 /cmd_vel은 ROS→Gazebo, odom·scan은 Gazebo→ROS 방향으로 타입을 변환한다.",
        ],
        "visual_capture": "captures/rviz_agv_sensors_actual.png",
        "visual_caption": "실제 RViz: URDF AGV body·wheel·camera·LiDAR·IMU frame",
        "capture_command": """printf 'Gazebo Sim version: '; gz sim --versions; printf '\\nros_gz_bridge prefix: '; ros2 pkg prefix ros_gz_bridge; printf '\\nSDF validation: '; gz sdf -k src/agv_gazebo/models/agv/model.sdf""",
        "checks": "Gazebo Harmonic 버전과 ros_gz_bridge 설치 경로가 출력되면 연동 도구가 준비된 상태다. 실제 주행 후 RViz에서 Fixed Frame=odom, Odometry=/odom, LaserScan=/scan을 설정한다.",
    },
    {
        "key": "D",
        "folder": "blocks/D_sensors",
        "title": "Block D — 센서",
        "modules": "M12 Camera · M13 2D LiDAR · M14 IMU · M15 운영/QoS",
        "goal": "Gazebo 센서를 생성하고 ROS topic, frame_id, /clock, update rate를 관리한다.",
        "files": [
            "agv_gazebo/models/agv/model.sdf",
            "agv_gazebo/config/bridge.yaml",
            "agv_sensors/agv_sensors/lidar_processor.py",
            "agv_sensors/agv_sensors/imu_monitor.py",
            "agv_bringup/config/sensors.yaml",
        ],
        "commands": """ros2 launch agv_bringup agv_sim.launch.py rviz:=false
# 새 터미널
ros2 topic hz /scan
ros2 topic echo /imu/data --once
ros2 run agv_sensors lidar_processor
ros2 topic echo /obstacle_distance
ros2 run rqt_image_view rqt_image_view""",
        "command_notes": [
            "SDF sensor의 update_rate·range·FOV가 Gazebo 원본 데이터를 만들고, bridge가 ROS 메시지로 변환한다.",
            "ros2 topic hz는 내용 대신 도착 주기를 측정한다. 현재 LiDAR 10 Hz, camera 15 Hz, IMU 100 Hz이다.",
            "front_half_angle_deg는 LiDAR processor가 장애물을 검사할 전방 반각을 정한다.",
        ],
        "capture_command": """printf 'SDF sensor declarations\\n'; rg -n '<sensor name=|<topic>|<update_rate>' src/agv_gazebo/models/agv/model.sdf; printf '\\nBridge mapping\\n'; cat src/agv_gazebo/config/bridge.yaml""",
        "checks": "SDF에 camera·lidar·imu sensor와 topic/update_rate가, bridge YAML에 ROS 메시지 타입과 방향이 함께 있어야 한다. 실행 중에는 topic hz와 header.frame_id를 반드시 확인한다.",
    },
    {
        "key": "E",
        "folder": "blocks/E_autonomy_logic",
        "title": "Block E — 제어·인지·미션",
        "modules": "M16 ros2_control · M17 YOLO · M18 장애물 인지 · M19 FSM · M20 PID",
        "goal": "센서·detection을 안전 우선순위와 상태머신으로 연결해 /cmd_vel을 만든다.",
        "files": [
            "agv_interfaces/msg/Detection.msg",
            "agv_interfaces/msg/DetectionArray.msg",
            "agv_vision/agv_vision/yolo_node.py",
            "agv_control/agv_control/safety_controller.py",
            "agv_mission/agv_mission/mission_manager.py",
        ],
        "commands": """cd /home/lab4090/ros2_curri/agv_ws
source install/setup.bash
ros2 interface show agv_interfaces/msg/Detection
ros2 run agv_sensors lidar_processor
ros2 run agv_control safety_controller --ros-args -p stop_distance:=0.5
ros2 run agv_mission mission_manager
ros2 topic echo /mission_state""",
        "command_notes": [
            "DetectionArray는 class·confidence·box 중심을 담고, yolo_node의 enable_yolo/model_path/threshold를 parameter로 조절한다.",
            "safety_controller는 /scan을 직접 검사해 stop_distance 안의 전진 명령을 zero Twist로 바꾼다.",
            "mission_manager는 10 Hz로 SEARCH·APPROACH·AVOID를 전환해 /cmd_vel_raw를 만든다.",
        ],
        "capture_command": """ros2 interface show agv_interfaces/msg/Detection; printf '\\n--- 실행 파일 ---\\n'; ros2 pkg executables agv_control; ros2 pkg executables agv_mission; ros2 pkg executables agv_vision""",
        "checks": "Detection 메시지 필드와 safety_controller·mission_manager 실행 파일이 출력되면 인지→안전→미션 구조가 빌드됐다. YOLO는 M17 문서의 가상환경 설치 후 enable_yolo=true로 활성화한다.",
    },
    {
        "key": "F",
        "folder": "blocks/F_integration",
        "title": "Block F — 통합·재현·최종 프로젝트",
        "modules": "M21 YAML/Launch/rosbag2 · M22 Final Project",
        "goal": "YAML과 하나의 launch 명령으로 8개 최종 AGV 패키지를 재현하고 rosbag으로 주행을 기록한다.",
        "files": [
            "agv_bringup/config/robot.yaml",
            "agv_bringup/config/sensors.yaml",
            "agv_bringup/config/vision.yaml",
            "agv_bringup/config/mission.yaml",
            "agv_bringup/launch/agv_sim.launch.py",
        ],
        "commands": """export PATH=/usr/bin:/bin:$PATH  # Miniconda 사용 시
source /opt/ros/jazzy/setup.bash
cd /home/lab4090/ros2_curri/agv_ws
colcon build --symlink-install --cmake-args -DPython3_EXECUTABLE=/usr/bin/python3
source install/setup.bash
ros2 launch agv_bringup agv_sim.launch.py
ros2 bag record -o ~/agv_bag_01 /scan /imu/data /odom /camera/image_raw""",
        "command_notes": [
            "colcon build는 8개 패키지를 의존성 순서로 install하고, source는 새 overlay를 현재 셸에 적용한다.",
            "agv_sim.launch.py는 Gazebo·bridge·TF·sensor·mission·safety를 한 명령으로 조립한다. rviz:=false로 GUI만 제외할 수 있다.",
            "ros2 bag record는 지정 토픽을 기록하며, --clock 재생은 use_sim_time 노드와 함께 사용한다.",
        ],
        "capture_command": """colcon list; printf '\\n--- 최종 launch 인자 ---\\n'; ros2 launch agv_bringup agv_sim.launch.py --show-args; printf '\\n--- 의존성 ---\\n'; rosdep check --from-paths src --ignore-src --rosdistro jazzy --skip-keys ament_python""",
        "checks": "8개 최종 AGV 패키지와 별도 C++ 교육 예제 1개, agv_sim.launch.py 인자가 출력되고 rosdep이 All system dependencies have been satisfied를 표시하면 재현 가능한 최종 구조다.",
    },
]


# 각 Block에서 처음 배우는 사람이 실제로 따라 만드는 파일 실습입니다.
# code는 슬라이드에 표시할 핵심 발췌이며, 전체 파일은 path에 있습니다.
WORKSHOPS = {
    "A": [
        {
            "title": "Python 노드 파일 만들기 (.py)",
            "path": "agv_ws/src/agv_control/agv_control/counter_publisher.py + setup.py",
            "make_command": "ros2 pkg create --build-type ament_python learning_py --dependencies rclpy std_msgs",
            "steps": [
                "1. ament_python 패키지를 만들고 agv_control/agv_control/에 .py 파일을 둡니다.",
                "2. Node·publisher·timer를 작성한 뒤 setup.py의 console_scripts에 실행 이름을 등록합니다.",
                "3. build → source → ros2 run 순서로 실행합니다. source가 없으면 새 실행 이름을 찾지 못합니다.",
            ],
            "code": "self.publisher = self.create_publisher(Int32, '/counter', 10)\nself.value = 0\nself.create_timer(1.0, self.publish_counter)\n\ndef publish_counter(self):\n    message = Int32(data=self.value)\n    self.publisher.publish(message)\n    self.value += 1\n\n# setup.py\n'counter_publisher = agv_control.counter_publisher:main'",
            "result": "실행하면 1초마다 /counter에 0, 1, 2…가 publish되고, counter_monitor가 같은 값을 받습니다.",
        },
        {
            "title": "C++ 노드 파일 만들기 (.cpp)",
            "path": "agv_ws/src/agv_cpp_examples/src/status_publisher.cpp + CMakeLists.txt",
            "make_command": "ros2 pkg create --build-type ament_cmake learning_cpp --dependencies rclcpp std_msgs",
            "steps": [
                "1. ros2 pkg create --build-type ament_cmake로 C++ 패키지와 src/ 폴더를 만듭니다.",
                "2. .cpp에서 rclcpp::Node, publisher, timer, spin을 작성합니다.",
                "3. CMakeLists.txt의 add_executable·ament_target_dependencies·install을 모두 작성합니다.",
            ],
            "code": "publisher_ = this->create_publisher<std_msgs::msg::String>(\n  \"/cpp_status\", 10);\ntimer_ = this->create_wall_timer(1s,\n  std::bind(&StatusPublisher::publish_status, this));\n\n# CMakeLists.txt\nadd_executable(status_publisher src/status_publisher.cpp)\nament_target_dependencies(status_publisher rclcpp std_msgs)\ninstall(TARGETS status_publisher\n  DESTINATION lib/${PROJECT_NAME})",
            "result": "ros2 run agv_cpp_examples status_publisher 뒤 /cpp_status에 'C++ AGV 상태 정상 #0'이 1초마다 나옵니다.",
        },
    ],
    "B": [
        {
            "title": "RViz용 로봇 조립 파일 만들기 (URDF/Xacro)",
            "path": "agv_ws/src/agv_description/urdf/agv.urdf.xacro",
            "make_command": "mkdir -p agv_description/urdf  # 패키지 안에 URDF/Xacro 폴더 준비",
            "steps": [
                "1. base_link에 visual·collision·inertial을 정의합니다.",
                "2. wheel_radius 같은 property를 먼저 선언하고 left/right wheel macro에 전달합니다.",
                "3. camera·LiDAR·IMU는 fixed joint로 base_link에 연결한 뒤 xacro/check_urdf로 검사합니다.",
            ],
            "code": "<xacro:property name='wheel_radius' value='0.08'/>\n<link name='base_link'>\n  <visual><geometry><box size='0.60 0.42 0.16'/></geometry></visual>\n</link>\n<xacro:agv_wheel side='left' y='${wheel_base/2}'\n                 radius='${wheel_radius}' width='${wheel_width}'/>\n<xacro:fixed_sensor name='lidar' parent='base_link'\n                    xyz='0.12 0 0.28' rpy='0 0 0' size='0.09 0.09 0.05'/>",
            "result": "xacro가 순수 URDF를 만들고, RViz에는 파란 body·두 바퀴·세 센서 frame이 보입니다.",
        },
        {
            "title": "Gazebo 물리·센서 파일 만들기 (SDF)",
            "path": "agv_ws/src/agv_gazebo/models/agv/model.sdf",
            "make_command": "mkdir -p agv_gazebo/models/agv  # model.config와 model.sdf를 같은 폴더에 둠",
            "steps": [
                "1. visual과 별도로 collision·inertial을 넣어 물리 엔진이 계산할 값을 만듭니다.",
                "2. LiDAR/camera/IMU sensor마다 topic과 update_rate를 정합니다.",
                "3. DiffDrive plugin의 joint 이름·wheel radius·separation을 URDF와 맞춥니다.",
            ],
            "code": "<sensor name='lidar' type='gpu_lidar'>\n  <topic>/scan</topic><update_rate>10</update_rate>\n  <lidar><scan><horizontal><samples>720</samples></horizontal></scan></lidar>\n</sensor>\n<plugin filename='gz-sim-diff-drive-system'>\n  <left_joint>left_wheel_joint</left_joint>\n  <right_joint>right_wheel_joint</right_joint>\n  <wheel_separation>0.38</wheel_separation><wheel_radius>0.08</wheel_radius>\n</plugin>",
            "result": "Gazebo는 충돌·질량을 계산하고 /scan·/imu/data·/camera/image_raw·/odom을 생성합니다.",
        },
    ],
    "C": [
        {
            "title": "Gazebo–ROS 연결 파일 만들기 (bridge.yaml)",
            "path": "agv_ws/src/agv_gazebo/config/bridge.yaml",
            "make_command": "mkdir -p agv_gazebo/config  # bridge.yaml을 launch와 같은 패키지에 둠",
            "steps": [
                "1. Gazebo의 실제 topic 이름을 gz topic -l로 확인합니다.",
                "2. ROS 이름·Gazebo 이름·양쪽 메시지 타입·방향을 한 항목에 적습니다.",
                "3. 명령은 ROS_TO_GZ, 센서와 odom은 GZ_TO_ROS인지 확인하고 launch를 다시 시작합니다.",
            ],
            "code": "- ros_topic_name: '/cmd_vel'\n  gz_topic_name: '/cmd_vel'\n  ros_type_name: 'geometry_msgs/msg/Twist'\n  gz_type_name: 'gz.msgs.Twist'\n  direction: ROS_TO_GZ\n- ros_topic_name: '/scan'\n  ros_type_name: 'sensor_msgs/msg/LaserScan'\n  gz_type_name: 'gz.msgs.LaserScan'\n  direction: GZ_TO_ROS",
            "result": "ROS의 /cmd_vel은 Gazebo DiffDrive로 가고, Gazebo /scan·/odom은 ROS에서 echo·RViz로 볼 수 있습니다.",
        },
    ],
    "D": [
        {
            "title": "LiDAR 처리 노드 만들기 (.py)",
            "path": "agv_ws/src/agv_sensors/agv_sensors/lidar_processor.py",
            "make_command": "ros2 pkg create --build-type ament_python agv_sensors --dependencies rclpy sensor_msgs std_msgs",
            "steps": [
                "1. LaserScan을 SensorDataQoS로 구독해 센서 publisher와 QoS를 맞춥니다.",
                "2. angle_min + index × angle_increment로 각 sample 방향을 계산합니다.",
                "3. 전방 ±15°의 유한 range 중 최솟값을 Float32로 publish합니다.",
            ],
            "code": "half_angle = radians(get_parameter('front_half_angle_deg').value)\nfront_ranges = [distance for index, distance in enumerate(scan.ranges)\n  if abs(scan.angle_min + index * scan.angle_increment) <= half_angle\n  and math.isfinite(distance)]\nresult = Float32(data=min(front_ranges) if front_ranges else float('inf'))\npublisher.publish(result)  # /obstacle_distance",
            "result": "전방 장애물이 0.48 m면 /obstacle_distance는 약 0.48, 측정값이 없으면 inf가 됩니다.",
        },
    ],
    "E": [
        {
            "title": "안전 필터 노드 만들기 (.py)",
            "path": "agv_ws/src/agv_control/agv_control/safety_controller.py",
            "make_command": "agv_control/agv_control/ 안에 safety_controller.py를 만들고 setup.py entry point를 추가",
            "steps": [
                "1. /scan과 /cmd_vel_raw를 구독하고 /cmd_vel publisher를 만듭니다.",
                "2. stop_distance와 front_half_angle_deg를 parameter로 선언합니다.",
                "3. 장애물이 가까운 전진 명령만 zero Twist로 바꿔 Gazebo에 전달합니다.",
            ],
            "code": "if obstacle_distance >= get_parameter('stop_distance').value\n   or command.linear.x <= 0.0:\n    safe = command\nelse:\n    safe = Twist()  # 모든 속도 0\n    logger.warn('stopping forward command')\npublisher.publish(safe)  # /cmd_vel",
            "result": "미션이 전진을 요청해도 LiDAR 전방 0.5 m 안에 물체가 있으면 /cmd_vel은 0이 됩니다.",
        },
    ],
    "F": [
        {
            "title": "숫자를 YAML 설정으로 분리하기 (.yaml)",
            "path": "agv_ws/src/agv_bringup/config/mission.yaml",
            "make_command": "mkdir -p agv_bringup/config  # node별 YAML을 config/에 저장",
            "steps": [
                "1. node 이름을 최상위 key로 쓰고 ros__parameters 아래에 값만 둡니다.",
                "2. 코드의 declare_parameter 이름과 YAML key 철자를 정확히 맞춥니다.",
                "3. 숫자를 바꾼 뒤 launch를 다시 시작해 코드 수정 없이 동작을 비교합니다.",
            ],
            "code": "mission_manager:\n  ros__parameters:\n    use_sim_time: true\n    stop_distance: 0.50\n    search_speed: 0.25\n    approach_speed: 0.15\n    image_center_x: 320\nsafety_controller:\n  ros__parameters:\n    stop_distance: 0.50",
            "result": "같은 0.50 m 정지 거리가 FSM과 safety node에 전달돼 서로 다른 기준으로 움직이지 않습니다.",
        },
        {
            "title": "여러 노드를 한 명령으로 실행하기 (launch.py)",
            "path": "agv_ws/src/agv_bringup/launch/agv_sim.launch.py",
            "make_command": "mkdir -p agv_bringup/launch  # launch.py는 패키지 setup.py의 data_files에도 등록",
            "steps": [
                "1. gazebo launch를 IncludeLaunchDescription으로 포함합니다.",
                "2. 각 Node에 해당 YAML 파일을 parameters로 전달합니다.",
                "3. rviz launch argument로 GUI를 켜거나 끄고, 실행 로그에서 각 process가 시작됐는지 확인합니다.",
            ],
            "code": "DeclareLaunchArgument('rviz', default_value='true')\nIncludeLaunchDescription(...gazebo.launch.py...)\nNode(package='agv_sensors', executable='lidar_processor',\n     parameters=[os.path.join(bringup, 'config', 'sensors.yaml')])\nNode(package='agv_mission', executable='mission_manager',\n     parameters=[os.path.join(bringup, 'config', 'mission.yaml')])\nNode(package='agv_control', executable='safety_controller',\n     parameters=[os.path.join(bringup, 'config', 'mission.yaml')])",
            "result": "ros2 launch 한 번으로 Gazebo·bridge·TF·센서·미션·안전 노드가 함께 시작됩니다.",
        },
    ],
}


def trim(text: str, max_lines: int = 20, max_width: int = 100) -> str:
    lines: list[str] = []
    for line in text.splitlines():
        lines.extend(textwrap.wrap(line, width=max_width, replace_whitespace=False) or [""])
    if len(lines) > max_lines:
        lines = lines[: max_lines - 1] + ["… (출력 일부 생략)"]
    return "\n".join(lines)


def terminal_capture(path: Path, title: str, command: str, output: str) -> None:
    width, height = 1920, 1080
    image = Image.new("RGB", (width, height), "#10151f")
    draw = ImageDraw.Draw(image)
    bold = ImageFont.truetype(FONT_BOLD_FILE, 32)
    regular = ImageFont.truetype(FONT_FILE, 20)
    small = ImageFont.truetype(FONT_FILE, 18)
    draw.rounded_rectangle((45, 35, width - 45, height - 35), radius=26, fill="#171f2e", outline="#34415b", width=2)
    draw.ellipse((75, 69, 93, 87), fill="#ff6258")
    draw.ellipse((104, 69, 122, 87), fill="#f5bf4f")
    draw.ellipse((133, 69, 151, 87), fill="#58c051")
    draw.text((180, 56), title + " — 실제 터미널 실행 캡처", font=bold, fill="#f2f5fb")
    draw.text((80, 130), "$ " + command.replace("\n", "\n$ "), font=small, fill="#8fd7ff", spacing=7)
    divider_y = 280
    draw.line((80, divider_y, width - 80, divider_y), fill="#41506e", width=2)
    draw.text((80, divider_y + 28), trim(output), font=regular, fill="#e6edf7", spacing=8)
    draw.text((80, height - 90), "캡처 생성 시각: " + datetime.now().strftime("%Y-%m-%d %H:%M:%S"), font=small, fill="#9aa9c1")
    image.save(path)


def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=RGBColor(33, 43, 61), bullet=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for index, line in enumerate(text.splitlines() or [""]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.font.name = FONT
        paragraph.font.size = Pt(max(size, 10))
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(5)
        if bullet and line:
            paragraph.level = 0
            paragraph.text = "• " + line
    return box


def add_header(slide, title: str, subtitle: str) -> None:
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.82))
    banner.fill.solid(); banner.fill.fore_color.rgb = RGBColor(23, 62, 111); banner.line.fill.background()
    add_textbox(slide, Inches(0.45), Inches(0.15), Inches(8.8), Inches(0.35), title, 25, True, RGBColor(255, 255, 255))
    add_textbox(slide, Inches(9.1), Inches(0.20), Inches(3.8), Inches(0.28), subtitle, 10, False, RGBColor(220, 234, 252))
    add_textbox(slide, Inches(0.45), Inches(7.17), Inches(12.3), Inches(0.18), "ROS 2 + Gazebo AGV 실습 · 모든 텍스트 최소 10pt", 10, False, RGBColor(102, 112, 128))


def add_bullets(slide, title: str, lines: list[str]) -> None:
    add_textbox(slide, Inches(0.65), Inches(1.05), Inches(12), Inches(0.45), title, 21, True, RGBColor(23, 62, 111))
    add_textbox(slide, Inches(0.8), Inches(1.65), Inches(11.7), Inches(4.9), "\n".join(lines), 15, False, bullet=True)


def add_code_slide(slide, block: dict) -> None:
    add_header(slide, block["title"], "명령어 실행 순서")
    add_textbox(slide, Inches(0.65), Inches(1.02), Inches(12), Inches(0.35), "권장 실행 순서", 21, True, RGBColor(23, 62, 111))
    code = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.55), Inches(12.0), Inches(4.95))
    code.fill.solid(); code.fill.fore_color.rgb = RGBColor(20, 28, 42); code.line.color.rgb = RGBColor(64, 80, 110)
    frame = code.text_frame; frame.clear(); frame.margin_left = Inches(0.26); frame.margin_top = Inches(0.20)
    for index, line in enumerate(block["commands"].splitlines()):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.font.name = FONT
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = RGBColor(223, 237, 255)
        paragraph.space_after = Pt(7)
    add_textbox(slide, Inches(0.75), Inches(6.66), Inches(11.8), Inches(0.28), "명령은 Block README의 순서와 동일합니다. 새 터미널에서도 ROS와 workspace를 source해야 합니다.", 10, False, RGBColor(80, 92, 112))


def add_command_notes_slide(slide, block: dict) -> None:
    add_header(slide, block["title"], "명령어 동작과 설정 포인트")
    add_bullets(slide, "이 명령이 실제로 하는 일", block["command_notes"])
    add_textbox(slide, Inches(0.8), Inches(5.95), Inches(11.7), Inches(0.55), "수치·토픽·방향을 바꿀 때는 해당 Block의 Mxx README에서 구현 파일과 기대 결과를 함께 확인합니다.", 12, False, RGBColor(63, 83, 115))


def add_workshop_slide(slide, block: dict, workshop: dict, number: int, total: int) -> None:
    add_header(slide, block["title"], f"초심자 파일 제작 실습 {number}/{total}")
    add_textbox(slide, Inches(0.55), Inches(0.92), Inches(12.1), Inches(0.30), workshop["title"], 19, True, RGBColor(23, 62, 111))
    add_textbox(slide, Inches(0.58), Inches(1.27), Inches(12.0), Inches(0.22), "실제 파일: " + workshop["path"], 11, False, RGBColor(76, 91, 116))
    add_textbox(slide, Inches(0.58), Inches(1.49), Inches(12.0), Inches(0.20), "처음 만들 때: " + workshop["make_command"], 10, False, RGBColor(23, 99, 140))

    code = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.78), Inches(6.35), Inches(4.83))
    code.fill.solid(); code.fill.fore_color.rgb = RGBColor(20, 28, 42); code.line.color.rgb = RGBColor(64, 80, 110)
    frame = code.text_frame; frame.clear(); frame.margin_left = Inches(0.20); frame.margin_top = Inches(0.15)
    for index, line in enumerate(workshop["code"].splitlines()):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.font.name = FONT
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = RGBColor(223, 237, 255)
        paragraph.space_after = Pt(3)

    add_textbox(slide, Inches(7.15), Inches(1.78), Inches(5.35), Inches(0.35), "만드는 순서", 16, True, RGBColor(23, 62, 111))
    add_textbox(slide, Inches(7.18), Inches(2.16), Inches(5.25), Inches(2.47), "\n".join(workshop["steps"]), 12, False, RGBColor(33, 43, 61), bullet=True)
    add_textbox(slide, Inches(7.18), Inches(4.90), Inches(5.25), Inches(0.35), "실행하면 이렇게 됩니다", 16, True, RGBColor(23, 62, 111))
    add_textbox(slide, Inches(7.18), Inches(5.35), Inches(5.25), Inches(0.90), workshop["result"], 12, False, RGBColor(33, 43, 61))


def add_visual_capture_slide(slide, block: dict, image_path: Path) -> None:
    add_header(slide, block["title"], "실제 Gazebo / RViz 화면")
    add_textbox(slide, Inches(0.75), Inches(0.98), Inches(11.9), Inches(0.30), block["visual_caption"], 16, True, RGBColor(23, 62, 111))
    slide.shapes.add_picture(str(image_path), Inches(3.35), Inches(1.35), height=Inches(5.50))


def create_presentation(block: dict, capture: Path) -> None:
    ppt = Presentation()
    ppt.slide_width = Inches(13.333)
    ppt.slide_height = Inches(7.5)
    blank = ppt.slide_layouts[6]

    cover = ppt.slides.add_slide(blank)
    cover.background.fill.solid(); cover.background.fill.fore_color.rgb = RGBColor(244, 248, 253)
    stripe = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = RGBColor(23, 62, 111); stripe.line.fill.background()
    add_textbox(cover, Inches(0.75), Inches(1.65), Inches(11.8), Inches(0.7), block["title"], 34, True, RGBColor(23, 62, 111))
    add_textbox(cover, Inches(0.78), Inches(2.52), Inches(11.6), Inches(0.4), block["modules"], 17, False, RGBColor(63, 83, 115))
    add_textbox(cover, Inches(0.78), Inches(3.4), Inches(11.4), Inches(0.8), block["goal"], 20, False, RGBColor(33, 43, 61))
    add_textbox(cover, Inches(0.78), Inches(6.45), Inches(11.4), Inches(0.35), "실제 실행 결과 캡처 포함 · 한국어 작성 · 최소 글꼴 크기 10pt", 12, False, RGBColor(89, 104, 126))

    overview = ppt.slides.add_slide(blank)
    add_header(overview, block["title"], "만들어진 파일")
    add_bullets(overview, "이 Block에서 실제로 만든 핵심 파일", block["files"])
    add_textbox(overview, Inches(0.8), Inches(5.85), Inches(11.7), Inches(0.55), "모듈별 만드는 방법과 상세 확인 기준은 이 폴더의 Mxx/README.md에 정리되어 있습니다.", 13, False, RGBColor(63, 83, 115))

    workshops = WORKSHOPS[block["key"]]
    for number, workshop in enumerate(workshops, start=1):
        workshop_slide = ppt.slides.add_slide(blank)
        add_workshop_slide(workshop_slide, block, workshop, number, len(workshops))

    commands = ppt.slides.add_slide(blank)
    add_code_slide(commands, block)

    command_notes = ppt.slides.add_slide(blank)
    add_command_notes_slide(command_notes, block)

    result = ppt.slides.add_slide(blank)
    add_header(result, block["title"], "실제 실행 결과")
    result.shapes.add_picture(str(capture), Inches(0.45), Inches(1.02), width=Inches(12.43), height=Inches(5.93))

    if "visual_capture" in block:
        visual = ROOT / block["folder"] / block["visual_capture"]
        visual_slide = ppt.slides.add_slide(blank)
        add_visual_capture_slide(visual_slide, block, visual)

    checklist = ppt.slides.add_slide(blank)
    add_header(checklist, block["title"], "확인 기준과 다음 단계")
    add_bullets(checklist, "통과 기준", [block["checks"], "문제가 생기면 README의 '확인' 명령부터 다시 실행하고, 첫 오류 메시지를 기준으로 수정합니다."])

    output = ROOT / block["folder"] / f"Block_{block['key']}_실습결과_명령어.pptx"
    ppt.save(output)


def main() -> None:
    for block in BLOCKS:
        block_dir = ROOT / block["folder"]
        capture_dir = block_dir / "captures"
        capture_dir.mkdir(exist_ok=True)
        output = shell(block["capture_command"])
        command_display = block["capture_command"].replace("; ", ";\n")
        capture = capture_dir / f"block_{block['key'].lower()}_actual_terminal.png"
        terminal_capture(capture, block["title"], command_display, output)
        create_presentation(block, capture)
        ppt_name = f"Block_{block['key']}_실습결과_명령어.pptx"
        print(f"created {capture.relative_to(ROOT)}")
        print(f"created {(block_dir / ppt_name).relative_to(ROOT)}")


if __name__ == "__main__":
    main()
